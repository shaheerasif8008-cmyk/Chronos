"""Bounded, non-executing previews for Chronos artifacts.

The preview contract is intentionally data-only. Office files, notebooks, and
archives are parsed into a bounded JSON representation; HTML/SVG is reduced to
a strict inert subset; React and other executable source is displayed as text.
No preview path evaluates macros, formulas, notebook cells, JavaScript, or
archive members.
"""
from __future__ import annotations

from html import escape
from html.parser import HTMLParser
from io import BytesIO
import json
import os
from pathlib import Path, PurePosixPath
import re
import subprocess
import sys
import tempfile
from typing import Any
from urllib.parse import urlparse
from zipfile import BadZipFile, ZipFile


class ArtifactPreviewError(ValueError):
    """Raised when an artifact cannot be safely previewed."""


_TEXT_LIMIT = 250_000
_MAX_BLOCKS = 500
_MAX_TABLE_ROWS = 200
_MAX_TABLE_COLS = 50
_MAX_SHEETS = 20
_MAX_SLIDES = 100
_MAX_NOTEBOOK_CELLS = 200
_MAX_ARCHIVE_ENTRIES = 500
_PREVIEW_TEXT_BUDGET = 1_000_000
_MAX_SHAPES_PER_SLIDE = 1_000
_MAX_TABLES = 25
_PDF_WORKER_TIMEOUT_SECONDS = 15
_PDF_WORKER_MAX_OUTPUT_BYTES = 24 * 1024 * 1024

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_IPYNB_MIMES = {"application/x-ipynb+json", "application/vnd.jupyter"}


class _PreviewBudget:
    """Cap attacker-controlled text copied into a JSON preview response."""

    def __init__(self, limit: int = _PREVIEW_TEXT_BUDGET) -> None:
        self.remaining = limit
        self.truncated = False

    @property
    def exhausted(self) -> bool:
        return self.remaining <= 0

    def take(self, value: Any, per_value_limit: int) -> str:
        text = str(value or "")
        allowed = min(per_value_limit, max(self.remaining, 0))
        if len(text) > allowed:
            self.truncated = True
        result = text[:allowed]
        self.remaining -= len(result)
        return result


def _record_budget_limit(result: dict[str, Any], budget: _PreviewBudget) -> None:
    if budget.truncated or budget.exhausted:
        result["limitations"].append(
            f"Preview text was truncated to {_PREVIEW_TEXT_BUDGET:,} characters."
        )


def _base(meta: dict[str, Any], content: bytes) -> dict[str, Any]:
    return {
        "status": "ready",
        "renderer": "download",
        "format": _format(meta),
        "mime_type": str(meta.get("mime_type") or "application/octet-stream"),
        "size_bytes": len(content),
        "limitations": [],
    }


def _format(meta: dict[str, Any]) -> str:
    mime = str(meta.get("mime_type") or "").split(";", 1)[0].strip().lower()
    title = str(meta.get("title") or "").lower()
    kind = str(meta.get("kind") or "").lower()
    suffix = PurePosixPath(title).suffix.lower()
    if mime == _DOCX_MIME or suffix == ".docx":
        return "docx"
    if mime == _PPTX_MIME or suffix == ".pptx":
        return "pptx"
    if mime == _XLSX_MIME or suffix in {".xlsx", ".xlsm"}:
        return "xlsx"
    if mime in _IPYNB_MIMES or suffix == ".ipynb" or kind == "notebook":
        return "ipynb"
    if mime == "application/pdf" or suffix == ".pdf":
        return "pdf"
    if mime in {"application/zip", "application/x-zip-compressed"} or suffix == ".zip" or kind == "project_bundle":
        return "zip"
    if mime in {"image/svg+xml", "image/svg"} or suffix == ".svg":
        return "svg"
    if mime in {"text/html", "application/xhtml+xml"} or suffix in {".html", ".htm"} or kind == "html":
        return "html"
    if kind in {"react", "tsx", "jsx"} or suffix in {".tsx", ".jsx"}:
        return "react"
    if kind in {"mermaid", "diagram"} or suffix in {".mmd", ".mermaid"}:
        return "diagram"
    if mime.startswith("image/"):
        return "image"
    if mime in {"application/json", "application/ld+json"} or suffix == ".json" or kind == "data":
        return "json"
    if mime in {"text/csv", "application/csv"} or suffix == ".csv" or kind == "csv":
        return "csv"
    if mime in {"text/markdown", "text/x-markdown"} or suffix in {".md", ".markdown"} or kind == "markdown":
        return "markdown"
    if mime.startswith("text/") or kind == "code":
        return "text"
    return "unknown"


def is_pdf_artifact(meta: dict[str, Any]) -> bool:
    """Use the same MIME/title classification for metadata and page routes."""
    return _format(meta) == "pdf"


def _decode(content: bytes) -> tuple[str, bool]:
    text = content.decode("utf-8", errors="replace")
    truncated = len(text) > _TEXT_LIMIT
    return text[:_TEXT_LIMIT], truncated


def _safe_zip(content: bytes, *, max_uncompressed_bytes: int) -> ZipFile:
    try:
        archive = ZipFile(BytesIO(content))
    except BadZipFile as exc:
        raise ArtifactPreviewError("The file is not a valid ZIP-based document.") from exc
    infos = archive.infolist()
    if len(infos) > _MAX_ARCHIVE_ENTRIES:
        archive.close()
        raise ArtifactPreviewError(f"Archive contains more than {_MAX_ARCHIVE_ENTRIES} entries.")
    total = 0
    for info in infos:
        name = info.filename.replace("\\", "/")
        path = PurePosixPath(name)
        if path.is_absolute() or ".." in path.parts:
            archive.close()
            raise ArtifactPreviewError("Archive contains an unsafe path.")
        if info.flag_bits & 0x1:
            archive.close()
            raise ArtifactPreviewError("Encrypted archives cannot be previewed safely.")
        total += int(info.file_size)
        if total > max_uncompressed_bytes:
            archive.close()
            raise ArtifactPreviewError("Archive expands beyond the configured preview limit.")
        if info.compress_size and info.file_size / max(info.compress_size, 1) > 1_000:
            archive.close()
            raise ArtifactPreviewError("Archive contains a suspicious compression ratio.")
    return archive


def build_preview(
    meta: dict[str, Any],
    content: bytes,
    *,
    max_bytes: int,
    max_uncompressed_bytes: int,
    max_pdf_pages: int,
) -> dict[str, Any]:
    """Return a bounded JSON-safe preview for an artifact."""
    result = _base(meta, content)
    if len(content) > max_bytes:
        result.update(status="unsupported", renderer="download")
        result["limitations"].append(
            f"Inline preview is limited to {max_bytes:,} bytes; download remains available."
        )
        return result

    fmt = result["format"]
    try:
        if fmt == "docx":
            return _preview_docx(result, content, max_uncompressed_bytes)
        if fmt == "pptx":
            return _preview_pptx(result, content, max_uncompressed_bytes)
        if fmt == "xlsx":
            return _preview_xlsx(result, content, max_uncompressed_bytes)
        if fmt == "ipynb":
            return _preview_notebook(result, content)
        if fmt == "zip":
            return _preview_zip(result, content, max_uncompressed_bytes)
        if fmt == "pdf":
            return _preview_pdf(result, content, max_pdf_pages)
        if fmt in {"html", "svg"}:
            return _preview_markup(result, content, svg=fmt == "svg")
        if fmt in {"react", "diagram"}:
            text, truncated = _decode(content)
            result.update(renderer="source", text=text)
            result["limitations"].append(
                "Source is shown without execution to protect the Chronos origin."
            )
            if truncated:
                result["limitations"].append("Preview text was truncated.")
            return result
        if fmt == "image":
            return _preview_image(result, content)
        if fmt in {"markdown", "text", "csv", "json"}:
            text, truncated = _decode(content)
            result.update(renderer=fmt, text=text)
            if truncated:
                result["limitations"].append("Preview text was truncated.")
            return result
    except ArtifactPreviewError as exc:
        result.update(status="error", renderer="download")
        result["limitations"].append(str(exc))
        return result
    except Exception:
        # Parser internals and document content must never leak through the API.
        result.update(status="error", renderer="download")
        result["limitations"].append(
            "Chronos could not safely parse this file. Download it to inspect it locally."
        )
        return result

    result.update(status="unsupported", renderer="download")
    result["limitations"].append(
        "This file type does not have a safe inline renderer. Download remains available."
    )
    return result


def _preview_docx(result: dict[str, Any], content: bytes, max_uncompressed: int) -> dict[str, Any]:
    from docx import Document

    with _safe_zip(content, max_uncompressed_bytes=max_uncompressed):
        pass
    document = Document(BytesIO(content))
    budget = _PreviewBudget()
    blocks: list[dict[str, Any]] = []
    for paragraph in document.paragraphs:
        raw_text = paragraph.text.strip()
        if not raw_text:
            continue
        text = budget.take(raw_text, 10_000)
        if not text:
            break
        style = str(paragraph.style.name or "") if paragraph.style else ""
        block_type = "heading" if style.lower().startswith("heading") else "paragraph"
        blocks.append({"type": block_type, "style": style[:300], "text": text})
        if len(blocks) >= _MAX_BLOCKS:
            result["limitations"].append("Document preview was truncated to 500 blocks.")
            break
    tables: list[list[list[str]]] = []
    for table in document.tables[:_MAX_TABLES]:
        if budget.exhausted:
            break
        rows: list[list[str]] = []
        for row in table.rows[:_MAX_TABLE_ROWS]:
            if budget.exhausted:
                break
            rows.append(
                [budget.take(cell.text, 2_000) for cell in row.cells[:_MAX_TABLE_COLS]]
            )
        tables.append(rows)
    if len(document.tables) > _MAX_TABLES:
        result["limitations"].append(f"Only the first {_MAX_TABLES} tables are shown.")
    image_count = len(document.inline_shapes)
    if image_count:
        result["limitations"].append(
            f"{image_count} embedded image{'s are' if image_count != 1 else ' is'} omitted from the safe text preview."
        )
    _record_budget_limit(result, budget)
    result.update(renderer="document", blocks=blocks, tables=tables)
    return result


def _preview_pptx(result: dict[str, Any], content: bytes, max_uncompressed: int) -> dict[str, Any]:
    from pptx import Presentation

    with _safe_zip(content, max_uncompressed_bytes=max_uncompressed):
        pass
    presentation = Presentation(BytesIO(content))
    budget = _PreviewBudget()
    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > _MAX_SLIDES:
            break
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        charts: list[dict[str, Any]] = []
        image_count = 0
        for shape_index, shape in enumerate(slide.shapes):
            if shape_index >= _MAX_SHAPES_PER_SLIDE:
                result["limitations"].append(
                    f"Slide {slide_index} was limited to {_MAX_SHAPES_PER_SLIDE:,} shapes."
                )
                break
            if budget.exhausted:
                break
            if getattr(shape, "has_text_frame", False):
                value = str(getattr(shape, "text", "")).strip()
                if value:
                    texts.append(budget.take(value, 10_000))
            if getattr(shape, "has_table", False) and len(tables) < _MAX_TABLES:
                rows: list[list[str]] = []
                for row in shape.table.rows[:_MAX_TABLE_ROWS]:
                    if budget.exhausted:
                        break
                    rows.append(
                        [budget.take(cell.text, 2_000) for cell in row.cells[:_MAX_TABLE_COLS]]
                    )
                tables.append(rows)
            if getattr(shape, "has_chart", False) and len(charts) < 25:
                chart = shape.chart
                series: list[dict[str, Any]] = []
                for item in list(chart.series)[:20]:
                    try:
                        values = [budget.take(value, 200) for value in list(item.values)[:100]]
                    except Exception:
                        values = []
                    series.append({"name": budget.take(getattr(item, "name", ""), 300), "values": values})
                charts.append({"series": series})
            if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
        slides.append({
            "number": slide_index,
            "texts": texts,
            "tables": tables,
            "charts": charts,
            "omitted_images": image_count,
        })
    if len(presentation.slides) > _MAX_SLIDES:
        result["limitations"].append(f"Only the first {_MAX_SLIDES} slides are shown.")
    if any(slide["omitted_images"] for slide in slides):
        result["limitations"].append("Slide images are omitted; text, tables, and cached chart data are shown.")
    _record_budget_limit(result, budget)
    result["limitations"].append("Animations, macros, media, and external links are never executed.")
    result.update(renderer="presentation", slides=slides)
    return result


def _cell_value(value: Any, budget: _PreviewBudget) -> str:
    if value is None:
        return ""
    if isinstance(value, (str, int, float, bool)):
        return budget.take(value, 2_000)
    return budget.take(value, 2_000)


def _preview_xlsx(result: dict[str, Any], content: bytes, max_uncompressed: int) -> dict[str, Any]:
    from openpyxl import load_workbook

    with _safe_zip(content, max_uncompressed_bytes=max_uncompressed):
        pass
    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True, keep_links=False)
    budget = _PreviewBudget()
    sheets: list[dict[str, Any]] = []
    try:
        for sheet in workbook.worksheets[:_MAX_SHEETS]:
            if budget.exhausted:
                break
            rows: list[list[str]] = []
            for row in sheet.iter_rows(max_row=_MAX_TABLE_ROWS, max_col=_MAX_TABLE_COLS, values_only=True):
                if budget.exhausted:
                    break
                rows.append([_cell_value(value, budget) for value in row])
            sheets.append({
                "name": budget.take(sheet.title, 300),
                "rows": rows,
                "truncated": sheet.max_row > _MAX_TABLE_ROWS or sheet.max_column > _MAX_TABLE_COLS,
            })
    finally:
        workbook.close()
    if len(workbook.sheetnames) > _MAX_SHEETS:
        result["limitations"].append(f"Only the first {_MAX_SHEETS} worksheets are shown.")
    if any(sheet["truncated"] for sheet in sheets):
        result["limitations"].append(
            f"Worksheets are limited to {_MAX_TABLE_ROWS} rows by {_MAX_TABLE_COLS} columns."
        )
    result["limitations"].append(
        "Formulas and macros are never executed; only stored cell values are shown."
    )
    _record_budget_limit(result, budget)
    result.update(renderer="workbook", sheets=sheets)
    return result


_ANSI = re.compile(r"\x1b\[[0-?]*[ -/]*[@-~]")


def _preview_notebook(result: dict[str, Any], content: bytes) -> dict[str, Any]:
    try:
        notebook = json.loads(content.decode("utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise ArtifactPreviewError("The notebook is not valid UTF-8 JSON.") from exc
    if not isinstance(notebook, dict) or not isinstance(notebook.get("cells"), list):
        raise ArtifactPreviewError("The notebook does not contain a valid cells array.")
    budget = _PreviewBudget()
    cells: list[dict[str, Any]] = []
    for index, raw_cell in enumerate(notebook["cells"][:_MAX_NOTEBOOK_CELLS], start=1):
        if budget.exhausted:
            break
        if not isinstance(raw_cell, dict):
            continue
        cell_type = budget.take(raw_cell.get("cell_type") or "raw", 50)
        source_raw = raw_cell.get("source") or ""
        source = "".join(source_raw) if isinstance(source_raw, list) else str(source_raw)
        bounded_source = budget.take(source, 20_000)
        outputs: list[str] = []
        if cell_type == "code":
            for output in (raw_cell.get("outputs") or [])[:25]:
                if not isinstance(output, dict):
                    continue
                text: Any = output.get("text")
                if text is None and isinstance(output.get("data"), dict):
                    text = output["data"].get("text/plain")
                if text is None and output.get("output_type") == "error":
                    text = f"{output.get('ename', 'Error')}: {output.get('evalue', '')}"
                if isinstance(text, list):
                    text = "".join(str(part) for part in text)
                if text is not None:
                    outputs.append(budget.take(_ANSI.sub("", str(text)), 10_000))
                if budget.exhausted:
                    break
        cells.append({
            "number": index,
            "cell_type": cell_type,
            "source": bounded_source,
            "outputs": outputs,
        })
    if len(notebook["cells"]) > _MAX_NOTEBOOK_CELLS:
        result["limitations"].append(f"Only the first {_MAX_NOTEBOOK_CELLS} cells are shown.")
    result["limitations"].append(
        "Notebook code is never executed. HTML, JavaScript, widgets, and rich stored outputs are omitted."
    )
    _record_budget_limit(result, budget)
    result.update(renderer="notebook", cells=cells)
    return result


def _preview_zip(result: dict[str, Any], content: bytes, max_uncompressed: int) -> dict[str, Any]:
    budget = _PreviewBudget()
    with _safe_zip(content, max_uncompressed_bytes=max_uncompressed) as archive:
        entries = [
            {
                "path": budget.take(info.filename, 2_000),
                "size_bytes": int(info.file_size),
                "compressed_bytes": int(info.compress_size),
                "directory": info.is_dir(),
            }
            for info in archive.infolist()
            if not budget.exhausted
        ]
    result["limitations"].append("Archive members are listed but never extracted or executed in Chronos.")
    _record_budget_limit(result, budget)
    result.update(renderer="archive", entries=entries)
    return result


def _preview_pdf(result: dict[str, Any], content: bytes, max_pages: int) -> dict[str, Any]:
    try:
        count = _run_pdf_worker("count", content, 0, max_pages)
        page_count = int(count.decode("ascii"))
    except Exception as exc:
        raise ArtifactPreviewError("The PDF could not be opened safely.") from exc
    result.update(renderer="pdf", page_count=page_count, preview_page_count=min(page_count, max_pages))
    if page_count > max_pages:
        result["limitations"].append(f"Inline preview is limited to the first {max_pages} pages.")
    result["limitations"].append("PDF pages are rendered to inert images; embedded actions are not executed.")
    return result


def render_pdf_page(content: bytes, page: int, *, max_pages: int) -> bytes:
    """Render a zero-based PDF page to PNG without exposing an active PDF viewer."""
    if page < 0 or page >= max_pages:
        raise ArtifactPreviewError("PDF preview page is outside the configured range.")
    try:
        png = _run_pdf_worker("render", content, page, max_pages)
        if not png.startswith(b"\x89PNG") or len(png) > _PDF_WORKER_MAX_OUTPUT_BYTES:
            raise ArtifactPreviewError("The PDF renderer returned an invalid image.")
        return png
    except ArtifactPreviewError:
        raise
    except Exception as exc:
        raise ArtifactPreviewError("The PDF page could not be rendered safely.") from exc


def _limit_preview_worker() -> None:
    """Best-effort POSIX resource envelope applied before native parsing."""
    try:
        import resource

        memory = 768 * 1024 * 1024
        limits = (
            (resource.RLIMIT_CPU, (10, 10)),
            (resource.RLIMIT_FSIZE, (32 * 1024 * 1024, 32 * 1024 * 1024)),
            (resource.RLIMIT_NOFILE, (32, 32)),
            (resource.RLIMIT_CORE, (0, 0)),
        )
        if hasattr(resource, "RLIMIT_AS"):
            limits += ((resource.RLIMIT_AS, (memory, memory)),)
        for limit, values in limits:
            try:
                resource.setrlimit(limit, values)
            except (OSError, ValueError):
                continue
    except Exception:
        pass


def _run_pdf_worker(action: str, content: bytes, page: int, max_pages: int) -> bytes:
    worker = Path(__file__).with_name("pdf_preview_worker.py")
    env = {
        "HOME": tempfile.gettempdir(),
        "PATH": os.defpath,
        "PYTHONHASHSEED": "random",
        "PYTHONIOENCODING": "utf-8",
        "PYTHONNOUSERSITE": "1",
    }
    try:
        completed = subprocess.run(
            [sys.executable, "-I", str(worker), action, str(page), str(max_pages)],
            input=content,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            cwd=tempfile.gettempdir(),
            env=env,
            timeout=_PDF_WORKER_TIMEOUT_SECONDS,
            check=False,
            preexec_fn=_limit_preview_worker if os.name == "posix" else None,
        )
    except subprocess.TimeoutExpired as exc:
        raise ArtifactPreviewError("PDF rendering exceeded the safe time limit.") from exc
    if completed.returncode != 0:
        message = completed.stderr.decode("utf-8", errors="replace").strip()
        raise ArtifactPreviewError(message or "PDF parser rejected the document.")
    if len(completed.stdout) > _PDF_WORKER_MAX_OUTPUT_BYTES:
        raise ArtifactPreviewError("PDF renderer exceeded the safe output limit.")
    return completed.stdout


def _preview_image(result: dict[str, Any], content: bytes) -> dict[str, Any]:
    worker = Path(__file__).with_name("image_preview_worker.py")
    env = {
        "HOME": tempfile.gettempdir(),
        "PATH": os.defpath,
        "PYTHONHASHSEED": "random",
        "PYTHONIOENCODING": "utf-8",
        "PYTHONNOUSERSITE": "1",
    }
    try:
        completed = subprocess.run(
            [sys.executable, "-I", str(worker)],
            input=content,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            cwd=tempfile.gettempdir(),
            env=env,
            timeout=10,
            check=False,
            preexec_fn=_limit_preview_worker if os.name == "posix" else None,
        )
    except subprocess.TimeoutExpired as exc:
        raise ArtifactPreviewError("Image validation exceeded the safe time limit.") from exc
    if completed.returncode != 0 or len(completed.stdout) > 4_096:
        message = completed.stderr.decode("utf-8", errors="replace").strip()
        raise ArtifactPreviewError(message or "Image parser rejected the document.")
    try:
        dimensions = json.loads(completed.stdout)
    except json.JSONDecodeError as exc:
        raise ArtifactPreviewError("Image validator returned an invalid result.") from exc
    image_format = str(dimensions["image_format"])
    declared = str(result.get("mime_type") or "").split(";", 1)[0].strip().lower()
    expected = {
        "image/png": "png",
        "image/jpeg": "jpeg",
        "image/jpg": "jpeg",
        "image/gif": "gif",
        "image/webp": "webp",
    }.get(declared)
    if expected is None or expected != image_format:
        raise ArtifactPreviewError("Image content does not match its declared browser format.")
    result.update(
        renderer="image",
        width=int(dimensions["width"]),
        height=int(dimensions["height"]),
        frames=int(dimensions["frames"]),
        image_format=image_format,
    )
    if result["frames"] > 1:
        result["limitations"].append(
            f"Animated preview is limited to {result['frames']} validated frames."
        )
    return result


_DROP_CONTENT = {"script", "style", "iframe", "frame", "object", "embed", "applet", "form", "template", "noscript"}
_HTML_TAGS = {
    "a", "abbr", "b", "blockquote", "br", "caption", "code", "col", "colgroup",
    "dd", "del", "details", "div", "dl", "dt", "em", "figcaption", "figure",
    "h1", "h2", "h3", "h4", "h5", "h6", "hr", "i", "img", "ins", "kbd",
    "li", "main", "mark", "ol", "p", "pre", "q", "s", "section", "small",
    "span", "strong", "sub", "summary", "sup", "table", "tbody", "td", "tfoot",
    "th", "thead", "tr", "u", "ul",
}
_SVG_TAGS = {
    "svg", "g", "path", "rect", "circle", "ellipse", "line", "polyline", "polygon",
    "text", "tspan", "defs", "lineargradient", "radialgradient", "stop", "clippath",
    "mask", "pattern", "title", "desc",
}
_GLOBAL_ATTRS = {"class", "title", "role", "aria-label", "aria-hidden"}
_HTML_ATTRS = {
    "a": {"href"}, "img": {"src", "alt", "width", "height"},
    "td": {"colspan", "rowspan"}, "th": {"colspan", "rowspan", "scope"},
    "col": {"span"}, "ol": {"start", "reversed", "type"}, "li": {"value"},
}
_SVG_ATTRS = {
    "viewbox", "width", "height", "x", "y", "x1", "x2", "y1", "y2", "cx", "cy",
    "r", "rx", "ry", "d", "points", "transform", "fill", "stroke", "stroke-width",
    "opacity", "fill-opacity", "stroke-opacity", "offset", "stop-color", "stop-opacity",
    "font-size", "font-family", "font-weight", "text-anchor", "preserveaspectratio",
}
_VOID = {"br", "hr", "img", "col"}


def _safe_data_image(value: str) -> bool:
    return bool(re.match(r"^data:image/(?:png|jpeg|gif|webp);base64,[A-Za-z0-9+/=\s]+$", value, re.I))


class _Sanitizer(HTMLParser):
    def __init__(self, *, svg: bool) -> None:
        super().__init__(convert_charrefs=True)
        self.allowed = _HTML_TAGS | (_SVG_TAGS if svg else set())
        self.parts: list[str] = []
        self.drop_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.lower()
        if tag in _DROP_CONTENT:
            self.drop_depth += 1
            return
        if self.drop_depth or tag not in self.allowed:
            return
        safe_attrs: list[str] = []
        for raw_name, raw_value in attrs:
            name = raw_name.lower()
            value = str(raw_value or "")[:2_000]
            if name.startswith("on") or name in {"style", "srcdoc", "xmlns:xlink", "xlink:href"}:
                continue
            allowed = name in _GLOBAL_ATTRS or name in _HTML_ATTRS.get(tag, set()) or name in _SVG_ATTRS
            if not allowed:
                continue
            if name == "href":
                parsed = urlparse(value)
                if parsed.scheme not in {"http", "https", "mailto"}:
                    continue
            if name == "src" and not _safe_data_image(value):
                continue
            if name in {"fill", "stroke"} and ("url(" in value.lower() or value.lower().startswith("data:")):
                continue
            safe_attrs.append(f'{name}="{escape(value, quote=True)}"')
        suffix = (" " + " ".join(safe_attrs)) if safe_attrs else ""
        self.parts.append(f"<{tag}{suffix}>")

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self.handle_starttag(tag, attrs)

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in _DROP_CONTENT:
            if self.drop_depth:
                self.drop_depth -= 1
            return
        if not self.drop_depth and tag in self.allowed and tag not in _VOID:
            self.parts.append(f"</{tag}>")

    def handle_data(self, data: str) -> None:
        if not self.drop_depth:
            self.parts.append(escape(data))


def _preview_markup(result: dict[str, Any], content: bytes, *, svg: bool) -> dict[str, Any]:
    text, truncated = _decode(content)
    sanitizer = _Sanitizer(svg=svg)
    sanitizer.feed(text)
    sanitizer.close()
    result.update(renderer="markup", html="".join(sanitizer.parts))
    result["limitations"].append(
        "Scripts, styles, forms, embeds, event handlers, and remote resources were removed."
    )
    if truncated:
        result["limitations"].append("Markup preview was truncated.")
    return result


def safe_download_headers(title: str | None, *, active_markup: bool = False) -> dict[str, str]:
    """Security and caching headers for authenticated/public artifact downloads."""
    candidate = re.sub(r"[^A-Za-z0-9._ -]+", "_", str(title or "artifact")).strip(" .")[:180]
    filename = candidate or "artifact"
    headers = {
        "Content-Disposition": f'attachment; filename="{filename}"',
        "X-Content-Type-Options": "nosniff",
        "Cache-Control": "private, no-store",
    }
    if active_markup:
        headers["Content-Security-Policy"] = "default-src 'none'; sandbox"
    return headers
