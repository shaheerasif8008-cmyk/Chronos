from __future__ import annotations

"""Document authoring connector — create and fill documents through the broker.

Routes from core.tool_broker for the ``doc.create``, ``doc.create_slides``,
``doc.fill_pdf`` and ``doc.render_chart`` tools. Like image_gen, this connector is
a pure renderer: it takes an explicit, structured spec and produces durable
artifacts. It contains no permissions, rate limiting, or planning — the agent
decides *what* to write and *where*; this module only renders it.

Design notes:

- ``doc.fill_pdf`` overlays a transparent reportlab layer onto the ORIGINAL pages
  (pypdf ``merge_page``). The source pixels are preserved exactly — this is the
  "fill the worksheet in place" path. Items place by explicit top-left point
  coordinates, or by locating an ``anchor_text`` string on the page and offsetting
  from it (so the agent can say "write the answer after this blank" without
  guessing pixel coordinates).
- All coordinates are PDF points (72 per inch) with a TOP-LEFT origin, which is
  more intuitive than reportlab's native bottom-left. Conversion happens here.
- Every output is saved via core.artifacts.save_artifact and returned as an
  artifact id, consistent with image_gen and the rest of the codebase.
- Failures degrade into an error ToolResult — they never propagate into the
  broker / SSE stream.
"""

import io
from typing import Any

from core.artifacts import get_artifact, read_artifact_content, save_artifact
from core.models import ToolResult

#: Default body font size for created documents.
_DEFAULT_FONT_SIZE = 11
#: Default overlay text size for fill_pdf.
_DEFAULT_OVERLAY_SIZE = 11


# ── helpers ──────────────────────────────────────────────────────────────────


def _normalize_ws(text: str) -> str:
    """Lowercase and collapse whitespace for tolerant substring matching."""
    return " ".join(str(text or "").split()).lower()


def _json_list_key(raw_json: str, key: str) -> list[dict[str, Any]]:
    """Parse model JSON and return the list of dict items under ``key`` (else [])."""
    import json

    try:
        parsed = json.loads(raw_json)
    except (json.JSONDecodeError, TypeError):
        return []
    if not isinstance(parsed, dict):
        return []
    value = parsed.get(key)
    if not isinstance(value, list):
        return []
    return [item for item in value if isinstance(item, dict)]


def _json_fields(raw_json: str) -> list[dict[str, Any]]:
    """Parse a vision response and return its ``fields`` list of dicts (else [])."""
    return _json_list_key(raw_json, "fields")


def _hex_to_rgb01(value: Any) -> tuple[float, float, float]:
    """Convert a ``#rrggbb`` (or ``rrggbb``) string to a 0..1 RGB triple.

    Falls back to black on any malformed input so rendering never raises on
    adversarial colour values.
    """
    s = str(value or "").lstrip("#")
    if len(s) != 6:
        return (0.0, 0.0, 0.0)
    try:
        r = int(s[0:2], 16) / 255.0
        g = int(s[2:4], 16) / 255.0
        b = int(s[4:6], 16) / 255.0
    except ValueError:
        return (0.0, 0.0, 0.0)
    return (r, g, b)


async def _resolve_pdf(artifact_id: str, org_id: str) -> bytes:
    """Resolve a PDF artifact to raw bytes, enforcing the org boundary."""
    meta = await get_artifact(artifact_id)
    if not meta or meta.get("is_deleted"):
        raise FileNotFoundError(f"artifact {artifact_id} not found")
    if str(meta.get("organization_id", "")) != str(org_id):
        raise PermissionError(f"artifact {artifact_id} does not belong to this organization")
    raw = await read_artifact_content(artifact_id) or b""
    if not raw:
        raise FileNotFoundError(f"artifact {artifact_id} has no content")
    return raw


async def _resolve_image(artifact_id: str, org_id: str) -> bytes | None:
    """Resolve an image artifact to bytes within the org, or None if unavailable."""
    meta = await get_artifact(artifact_id)
    if not meta or meta.get("is_deleted"):
        return None
    if str(meta.get("organization_id", "")) != str(org_id):
        return None
    return await read_artifact_content(artifact_id)


def _anchor_positions(raw: bytes) -> dict[int, list[tuple[str, float, float]]]:
    """Map 0-based page index → list of (text, x, y_bottom_left) word positions.

    Uses pypdf's text-extraction visitor to capture the device position of each
    text fragment. Coordinates are in PDF points with reportlab's native
    bottom-left origin. Best-effort: returns {} if extraction fails.
    """
    from pypdf import PdfReader

    positions: dict[int, list[tuple[str, float, float]]] = {}
    try:
        reader = PdfReader(io.BytesIO(raw))
    except Exception:
        return positions

    for idx, page in enumerate(reader.pages):
        found: list[tuple[str, float, float]] = []

        def visitor(text, cm, tm, font_dict, font_size, _found=found):  # noqa: ANN001
            t = (text or "").strip()
            if t:
                _found.append((t, float(tm[4]), float(tm[5])))

        try:
            page.extract_text(visitor_text=visitor)
        except Exception:
            pass
        positions[idx] = found
    return positions


def _find_anchor(
    page_words: list[tuple[str, float, float]], anchor_text: str
) -> tuple[float, float] | None:
    """Return the (x, y_bottom_left) of the first word fragment containing anchor_text."""
    needle = anchor_text.strip().lower()
    if not needle:
        return None
    for text, x, y in page_words:
        if needle in text.lower():
            return (x, y)
    return None


# ── connector ────────────────────────────────────────────────────────────────


class DocAuthoringConnector:
    async def execute(self, tool: str, args: dict[str, Any]) -> ToolResult:
        args.pop("__connector_tier", None)
        org_id = str(args.pop("__org_id", "default") or "default")
        task_id = args.pop("__task_id", None)

        if tool == "doc.create":
            return await self._create(args, org_id=org_id, task_id=task_id)
        if tool == "doc.create_slides":
            return await self._create_slides(args, org_id=org_id, task_id=task_id)
        if tool == "doc.fill_pdf":
            return await self._fill_pdf(args, org_id=org_id, task_id=task_id)
        if tool == "doc.render_chart":
            return await self._render_chart(args, org_id=org_id, task_id=task_id)
        if tool == "doc.detect_fields":
            return await self._detect_fields(args, org_id=org_id)
        if tool == "doc.verify_fill":
            return await self._verify_fill(args, org_id=org_id)
        raise ValueError(f"Unknown doc authoring tool: {tool}")

    # ── doc.create ────────────────────────────────────────────────────────────

    async def _create(
        self, args: dict[str, Any], *, org_id: str, task_id: str | None
    ) -> ToolResult:
        """Author a new PDF / DOCX / Markdown document from structured blocks.

        blocks: list of {type, ...} where type is one of:
          - "heading"   {text, level?}      level 1..3
          - "paragraph" {text}
          - "bullet"    {text}
          - "image"     {artifact_id, width?, height?}
          - "pagebreak" {}
        """
        fmt = str(args.get("format") or "pdf").lower()
        title = str(args.get("title") or "Document")
        blocks = args.get("blocks")
        if not isinstance(blocks, list) or not blocks:
            return ToolResult(
                data={"status": "error", "reason": "blocks (non-empty list) is required"},
                summary="doc.create: blocks is required",
            )

        try:
            if fmt == "markdown" or fmt == "md":
                raw, mime, kind = self._render_markdown(title, blocks), "text/markdown", "markdown"
            elif fmt == "docx":
                raw, mime, kind = await self._render_docx(title, blocks, org_id), (
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ), "file"
            elif fmt == "pdf":
                raw, mime, kind = await self._render_pdf(title, blocks, org_id), "application/pdf", "file"
            else:
                return ToolResult(
                    data={"status": "error", "reason": f"unsupported format: {fmt}"},
                    summary=f"doc.create: unsupported format {fmt!r} (use pdf, docx, or markdown)",
                )
        except Exception as exc:
            return ToolResult(
                data={"status": "error", "reason": f"render failed: {type(exc).__name__}: {exc}"},
                summary=f"doc.create failed: {exc}",
            )

        artifact_id = await save_artifact(
            raw, kind=kind, title=title, task_id=task_id, org_id=org_id,
            mime_type=mime, created_by="doc_authoring_connector",
        )
        return ToolResult(
            data={"status": "success", "artifact_id": artifact_id, "format": fmt, "blocks": len(blocks)},
            summary=f"Created {fmt.upper()} '{title}' ({len(blocks)} blocks) → artifact {artifact_id}",
        )

    def _render_markdown(self, title: str, blocks: list[dict[str, Any]]) -> bytes:
        lines: list[str] = [f"# {title}", ""]
        for b in blocks:
            t = str(b.get("type") or "paragraph")
            text = str(b.get("text") or "")
            if t == "heading":
                level = max(1, min(3, int(b.get("level") or 2)))
                lines.append(f"{'#' * (level + 1)} {text}")
            elif t == "bullet":
                lines.append(f"- {text}")
            elif t == "pagebreak":
                lines.append("\n---\n")
            elif t == "image":
                lines.append(f"![image](artifact:{b.get('artifact_id', '')})")
            else:
                lines.append(text)
            lines.append("")
        return "\n".join(lines).encode("utf-8")

    async def _render_docx(self, title: str, blocks: list[dict[str, Any]], org_id: str) -> bytes:
        from docx import Document
        from docx.shared import Inches

        doc = Document()
        doc.add_heading(title, level=0)
        for b in blocks:
            t = str(b.get("type") or "paragraph")
            text = str(b.get("text") or "")
            if t == "heading":
                doc.add_heading(text, level=max(1, min(3, int(b.get("level") or 2))))
            elif t == "bullet":
                doc.add_paragraph(text, style="List Bullet")
            elif t == "pagebreak":
                doc.add_page_break()
            elif t == "image":
                img = await _resolve_image(str(b.get("artifact_id") or ""), org_id)
                if img:
                    width = b.get("width")
                    doc.add_picture(io.BytesIO(img), width=Inches(float(width)) if width else None)
            else:
                doc.add_paragraph(text)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    async def _render_pdf(self, title: str, blocks: list[dict[str, Any]], org_id: str) -> bytes:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            Image as RLImage,
            ListFlowable,
            ListItem,
            PageBreak,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
        )

        styles = getSampleStyleSheet()
        story: list[Any] = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
        for b in blocks:
            t = str(b.get("type") or "paragraph")
            text = str(b.get("text") or "")
            if t == "heading":
                level = max(1, min(3, int(b.get("level") or 2)))
                story.append(Paragraph(text, styles[f"Heading{level}"]))
            elif t == "bullet":
                story.append(ListFlowable([ListItem(Paragraph(text, styles["BodyText"]))], bulletType="bullet"))
            elif t == "pagebreak":
                story.append(PageBreak())
            elif t == "image":
                img = await _resolve_image(str(b.get("artifact_id") or ""), org_id)
                if img:
                    width = float(b.get("width") or 4) * inch
                    height = float(b.get("height") or 3) * inch
                    story.append(RLImage(io.BytesIO(img), width=width, height=height))
            else:
                story.append(Paragraph(text, styles["BodyText"]))
            story.append(Spacer(1, 8))

        buf = io.BytesIO()
        SimpleDocTemplate(buf, pagesize=letter).build(story)
        return buf.getvalue()

    # ── doc.create_slides ─────────────────────────────────────────────────────

    async def _create_slides(
        self, args: dict[str, Any], *, org_id: str, task_id: str | None
    ) -> ToolResult:
        """Author a PPTX deck from a slides list.

        slides: list of {title?, bullets?: [str], image_artifact_id?, notes?}
        """
        slides = args.get("slides")
        title = str(args.get("title") or "Presentation")
        if not isinstance(slides, list) or not slides:
            return ToolResult(
                data={"status": "error", "reason": "slides (non-empty list) is required"},
                summary="doc.create_slides: slides is required",
            )
        try:
            raw = await self._render_pptx(slides, org_id)
        except Exception as exc:
            return ToolResult(
                data={"status": "error", "reason": f"render failed: {type(exc).__name__}: {exc}"},
                summary=f"doc.create_slides failed: {exc}",
            )
        artifact_id = await save_artifact(
            raw, kind="file", title=title, task_id=task_id, org_id=org_id,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            created_by="doc_authoring_connector",
        )
        return ToolResult(
            data={"status": "success", "artifact_id": artifact_id, "slides": len(slides)},
            summary=f"Created PPTX '{title}' ({len(slides)} slides) → artifact {artifact_id}",
        )

    async def _render_pptx(self, slides: list[dict[str, Any]], org_id: str) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        for s in slides:
            slide_title = str(s.get("title") or "")
            bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
            image_id = str(s.get("image_artifact_id") or "")

            layout = prs.slide_layouts[1] if bullets else prs.slide_layouts[5]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = slide_title

            if bullets:
                body = slide.placeholders[1].text_frame
                body.clear()
                for i, line in enumerate(bullets):
                    para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    para.text = str(line)
                    para.font.size = Pt(18)

            if image_id:
                img = await _resolve_image(image_id, org_id)
                if img:
                    slide.shapes.add_picture(
                        io.BytesIO(img), Inches(5.5), Inches(1.5), width=Inches(4)
                    )

            notes = s.get("notes")
            if notes:
                slide.notes_slide.notes_text_frame.text = str(notes)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── doc.fill_pdf ──────────────────────────────────────────────────────────

    async def _fill_pdf(
        self, args: dict[str, Any], *, org_id: str, task_id: str | None
    ) -> ToolResult:
        """Overlay text/images onto an existing PDF, preserving the original pages.

        items: list of {page, type, ...} (page is 1-based):
          text item:  {page, type:"text", text, x?, y?, anchor_text?, dx?, dy?,
                       size?, color?, font?}
          image item: {page, type:"image", artifact_id, width, height,
                       x?, y?, anchor_text?, dx?, dy?}

        Coordinates x/y are PDF points, TOP-LEFT origin. With anchor_text, the
        item is placed at the anchor's location offset by (dx, dy) points
        (dx>0 right, dy>0 down).
        """
        from pypdf import PdfReader, PdfWriter
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas

        source_id = str(args.get("artifact_id") or "")
        items = args.get("items")
        if not source_id:
            return ToolResult(
                data={"status": "error", "reason": "artifact_id is required"},
                summary="doc.fill_pdf: artifact_id is required",
            )
        if not isinstance(items, list) or not items:
            return ToolResult(
                data={"status": "error", "reason": "items (non-empty list) is required"},
                summary="doc.fill_pdf: items is required",
            )

        try:
            raw = await _resolve_pdf(source_id, org_id)
        except (FileNotFoundError, PermissionError) as exc:
            return ToolResult(
                data={"status": "error", "reason": str(exc)},
                summary=f"doc.fill_pdf: {exc}",
            )

        try:
            reader = PdfReader(io.BytesIO(raw))
        except Exception as exc:
            return ToolResult(
                data={"status": "error", "reason": f"could not read PDF: {exc}"},
                summary=f"doc.fill_pdf: source artifact is not a readable PDF ({exc})",
            )

        page_count = len(reader.pages)
        anchors = _anchor_positions(raw)

        # Group items by 0-based page index, dropping out-of-range pages honestly.
        by_page: dict[int, list[dict[str, Any]]] = {}
        skipped: list[str] = []
        for item in items:
            if not isinstance(item, dict):
                continue
            pidx = int(item.get("page") or 1) - 1
            if pidx < 0 or pidx >= page_count:
                skipped.append(f"page {pidx + 1} out of range (1..{page_count})")
                continue
            by_page.setdefault(pidx, []).append(item)

        writer = PdfWriter(clone_from=reader)
        placed = 0
        for pidx, page in enumerate(writer.pages):
            page_items = by_page.get(pidx)
            if not page_items:
                continue
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            overlay_buf = io.BytesIO()
            c = canvas.Canvas(overlay_buf, pagesize=(width, height))
            page_placed = 0
            for item in page_items:
                drew = await self._draw_item(
                    c, item, page_height=height, page_words=anchors.get(pidx, []),
                    org_id=org_id, image_reader=ImageReader,
                )
                if drew:
                    page_placed += 1
            # An overlay with nothing drawn (e.g. all anchors missed) yields a
            # zero-page reportlab PDF — skip the merge entirely in that case.
            if page_placed == 0:
                continue
            c.save()
            overlay_buf.seek(0)
            page.merge_page(PdfReader(overlay_buf).pages[0])
            placed += page_placed

        out = io.BytesIO()
        writer.write(out)
        filled_bytes = out.getvalue()

        artifact_id = await save_artifact(
            filled_bytes, kind="file",
            title=f"Filled: {args.get('title') or source_id}",
            task_id=task_id, org_id=org_id, mime_type="application/pdf",
            parent_artifact_id=source_id, created_by="doc_authoring_connector",
        )
        return ToolResult(
            data={
                "status": "success",
                "artifact_id": artifact_id,
                "source_artifact_id": source_id,
                "items_placed": placed,
                "items_skipped": skipped,
                "page_count": page_count,
            },
            summary=(
                f"Filled PDF: placed {placed} item(s) over {page_count} page(s) "
                f"→ artifact {artifact_id}"
                + (f"; skipped {len(skipped)}" if skipped else "")
            ),
        )

    async def _draw_item(
        self, c, item: dict[str, Any], *, page_height: float,
        page_words: list[tuple[str, float, float]], org_id: str, image_reader,
    ) -> bool:
        """Draw a single overlay item onto the reportlab canvas. Returns True if drawn."""
        itype = str(item.get("type") or "text")

        # Resolve placement. reportlab uses a bottom-left origin; user coords are
        # top-left. y_tl is the top-left y in points; convert per-item below.
        anchor_text = item.get("anchor_text")
        if anchor_text:
            hit = _find_anchor(page_words, str(anchor_text))
            if hit is None:
                return False
            ax, ay_bl = hit  # anchor x, y (bottom-left origin)
            x = ax + float(item.get("dx") or 0)
            # anchor top-left y, then offset down by dy
            y_tl = (page_height - ay_bl) + float(item.get("dy") or 0)
        else:
            x = float(item.get("x") or 0)
            y_tl = float(item.get("y") or 0)

        if itype == "text":
            text = str(item.get("text") or "")
            if not text:
                return False
            size = float(item.get("size") or _DEFAULT_OVERLAY_SIZE)
            font = str(item.get("font") or "Helvetica")
            c.setFillColorRGB(*_hex_to_rgb01(item.get("color")))
            try:
                c.setFont(font, size)
            except Exception:
                c.setFont("Helvetica", size)
            # Baseline sits ~size below the top-left y.
            c.drawString(x, page_height - y_tl - size, text)
            return True

        if itype == "image":
            img = await _resolve_image(str(item.get("artifact_id") or ""), org_id)
            if not img:
                return False
            w = float(item.get("width") or 100)
            h = float(item.get("height") or 100)
            try:
                c.drawImage(
                    image_reader(io.BytesIO(img)),
                    x, page_height - y_tl - h, width=w, height=h, mask="auto",
                )
            except Exception:
                return False
            return True

        return False

    # ── doc.render_chart ──────────────────────────────────────────────────────

    async def _render_chart(
        self, args: dict[str, Any], *, org_id: str, task_id: str | None
    ) -> ToolResult:
        """Render a precise, code-generated chart via matplotlib → image artifact.

        spec: {chart_type: line|bar|scatter|pie, series|values, labels?, title?,
               xlabel?, ylabel?}
        """
        chart_type = str(args.get("chart_type") or "line").lower()
        title = str(args.get("title") or "Chart")
        try:
            raw = self._render_chart_png(chart_type, args)
        except Exception as exc:
            return ToolResult(
                data={"status": "error", "reason": f"render failed: {type(exc).__name__}: {exc}"},
                summary=f"doc.render_chart failed: {exc}",
            )
        artifact_id = await save_artifact(
            raw, kind="image", title=f"Chart: {title}", task_id=task_id, org_id=org_id,
            mime_type="image/png", created_by="doc_authoring_connector",
        )
        return ToolResult(
            data={"status": "success", "artifact_id": artifact_id, "chart_type": chart_type},
            summary=f"Rendered {chart_type} chart '{title}' → artifact {artifact_id}",
        )

    def _render_chart_png(self, chart_type: str, args: dict[str, Any]) -> bytes:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(6, 4))
        labels = args.get("labels") if isinstance(args.get("labels"), list) else None

        if chart_type == "pie":
            values = [float(v) for v in (args.get("values") or [])]
            ax.pie(values, labels=labels, autopct="%1.1f%%")
        else:
            series = args.get("series")
            if not isinstance(series, list):
                # Single implicit series from "values".
                series = [{"name": "", "values": args.get("values") or []}]
            for s in series:
                ys = [float(v) for v in (s.get("values") or [])]
                xs = labels if (labels and len(labels) == len(ys)) else list(range(len(ys)))
                name = str(s.get("name") or "")
                if chart_type == "bar":
                    ax.bar(xs, ys, label=name)
                elif chart_type == "scatter":
                    ax.scatter(xs, ys, label=name)
                else:
                    ax.plot(xs, ys, marker="o", label=name)
            if any(str(s.get("name") or "") for s in series):
                ax.legend()
            if args.get("xlabel"):
                ax.set_xlabel(str(args["xlabel"]))
            if args.get("ylabel"):
                ax.set_ylabel(str(args["ylabel"]))

        ax.set_title(str(args.get("title") or "Chart"))
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=120)
        plt.close(fig)
        return buf.getvalue()

    # ── doc.detect_fields ─────────────────────────────────────────────────────

    async def _detect_fields(self, args: dict[str, Any], *, org_id: str) -> ToolResult:
        """Detect fillable regions in a PDF via vision, with absolute coordinates.

        Renders each requested page to an image, asks the vision model for the
        fillable regions (blanks, answer lines, checkboxes) as normalized boxes,
        and converts them to PDF points (top-left origin). The returned ``fields``
        carry a ``fill_hint`` ready to drop straight into ``doc.fill_pdf`` items —
        so the agent never has to guess coordinates.
        """
        from core.llm import vision_json
        from pypdf import PdfReader

        source_id = str(args.get("artifact_id") or "")
        if not source_id:
            return ToolResult(
                data={"status": "error", "reason": "artifact_id is required"},
                summary="doc.detect_fields: artifact_id is required",
            )
        try:
            raw = await _resolve_pdf(source_id, org_id)
        except (FileNotFoundError, PermissionError) as exc:
            return ToolResult(
                data={"status": "error", "reason": str(exc)},
                summary=f"doc.detect_fields: {exc}",
            )
        try:
            reader = PdfReader(io.BytesIO(raw))
        except Exception as exc:
            return ToolResult(
                data={"status": "error", "reason": f"could not read PDF: {exc}"},
                summary=f"doc.detect_fields: source is not a readable PDF ({exc})",
            )

        page_count = len(reader.pages)
        requested = args.get("pages")
        if isinstance(requested, list) and requested:
            page_indices = [int(p) - 1 for p in requested if 0 < int(p) <= page_count]
        else:
            page_indices = list(range(page_count))

        instruction = (
            "You are analysing a form/worksheet page to find every place a person "
            "should write an answer (blank lines, empty boxes, answer fields, "
            "checkboxes). Return ONLY JSON of this exact shape:\n"
            '{"fields": [{"label": "the question or field label", '
            '"field_type": "text"|"checkbox", '
            '"bbox": [x0, y0, x1, y1]}]}\n'
            "Coordinates are fractions from 0.0 to 1.0 of the page, with the ORIGIN "
            "at the TOP-LEFT (x to the right, y downward). bbox is the empty region "
            "to be filled, NOT the printed label. Return [] if there are no fields."
        )

        all_fields: list[dict[str, Any]] = []
        pages_analyzed = 0
        for pidx in page_indices:
            page = reader.pages[pidx]
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            png = self._render_page_png(raw, pidx)
            if not png:
                continue
            content = await vision_json(png, "image/png", instruction)
            if not content:
                continue
            pages_analyzed += 1
            for field in _json_fields(content):
                bbox = field.get("bbox")
                if not isinstance(bbox, list) or len(bbox) != 4:
                    continue
                try:
                    x0, y0, x1, y1 = (max(0.0, min(1.0, float(v))) for v in bbox)
                except (TypeError, ValueError):
                    continue
                px, py = x0 * width, y0 * height
                pw, ph = max(0.0, (x1 - x0)) * width, max(0.0, (y1 - y0)) * height
                all_fields.append({
                    "page": pidx + 1,
                    "label": str(field.get("label") or ""),
                    "field_type": str(field.get("field_type") or "text"),
                    "bbox_points": {"x": round(px, 1), "y": round(py, 1),
                                    "width": round(pw, 1), "height": round(ph, 1)},
                    # Ready to drop into doc.fill_pdf: place text just inside the box.
                    "fill_hint": {"x": round(px + 2, 1),
                                  "y": round(py + ph * 0.25, 1),
                                  "size": round(min(14.0, max(8.0, ph * 0.6)), 1)},
                })

        if pages_analyzed == 0:
            return ToolResult(
                data={"status": "unavailable", "fields": [],
                      "fallback_reason": "no vision model configured or rendering failed",
                      "page_count": page_count},
                summary=("doc.detect_fields: vision analysis unavailable — fall back to "
                         "anchor_text placement in doc.fill_pdf."),
            )
        return ToolResult(
            data={"status": "success", "fields": all_fields,
                  "pages_analyzed": pages_analyzed, "page_count": page_count},
            summary=f"Detected {len(all_fields)} fillable field(s) across {pages_analyzed} page(s).",
        )

    def _render_page_png(self, raw: bytes, page_index: int) -> bytes | None:
        """Render one PDF page to PNG bytes via pypdfium2 (no system binary)."""
        try:
            import pypdfium2 as pdfium

            pdf = pdfium.PdfDocument(raw)
            try:
                pil = pdf[page_index].render(scale=2.0).to_pil()
            finally:
                pdf.close()
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            return None

    # ── doc.verify_fill ───────────────────────────────────────────────────────

    async def _verify_fill(self, args: dict[str, Any], *, org_id: str) -> ToolResult:
        """Re-read a filled PDF and verify each expected answer landed legibly.

        Re-parses the filled PDF (text extraction, OCR fallback) and checks that
        each expected answer string is present in the rendered output. Optionally
        runs an LLM correctness re-check against the supplied questions. This is
        the "recheck its work" step after doc.fill_pdf.
        """
        from parsing.engine import parse_document

        source_id = str(args.get("artifact_id") or "")
        expected = args.get("expected")
        if not source_id:
            return ToolResult(
                data={"status": "error", "reason": "artifact_id is required"},
                summary="doc.verify_fill: artifact_id is required",
            )
        if not isinstance(expected, list) or not expected:
            return ToolResult(
                data={"status": "error", "reason": "expected (non-empty list) is required"},
                summary="doc.verify_fill: expected answers list is required",
            )
        try:
            raw = await _resolve_pdf(source_id, org_id)
        except (FileNotFoundError, PermissionError) as exc:
            return ToolResult(
                data={"status": "error", "reason": str(exc)},
                summary=f"doc.verify_fill: {exc}",
            )

        doc = await parse_document(raw, "application/pdf", "filled.pdf")
        haystack = _normalize_ws(doc.full_text)

        items: list[dict[str, Any]] = []
        missing: list[str] = []
        for entry in expected:
            if isinstance(entry, dict):
                answer = str(entry.get("answer") or "")
                label = str(entry.get("label") or "")
            else:
                answer, label = str(entry), ""
            present = bool(answer) and _normalize_ws(answer) in haystack
            items.append({"label": label, "answer": answer, "present": present})
            if not present:
                missing.append(answer or label)

        result: dict[str, Any] = {
            "status": "success",
            "all_present": not missing,
            "items": items,
            "missing": missing,
            "parser_used": doc.parser_used,
        }

        # Optional correctness re-check (the "recheck its work" pass).
        if args.get("recheck_correctness"):
            result["correctness"] = await self._recheck_correctness(doc.full_text, expected)

        placed = sum(1 for it in items if it["present"])
        return ToolResult(
            data=result,
            summary=(f"Verified filled PDF: {placed}/{len(items)} expected answer(s) present"
                     + (f"; missing {len(missing)}" if missing else "")),
        )

    async def _recheck_correctness(
        self, filled_text: str, expected: list[Any]
    ) -> dict[str, Any]:
        """LLM re-check of answer correctness against the filled content. Degrades honestly."""
        from core import llm

        qa_lines = []
        for entry in expected:
            if isinstance(entry, dict) and entry.get("question"):
                qa_lines.append(f"Q: {entry.get('question')}\nA: {entry.get('answer')}")
        if not qa_lines:
            return {"status": "skipped", "reason": "no questions supplied for correctness check"}

        prompt = (
            "You are checking a completed worksheet. For each question/answer pair, "
            "decide if the answer is correct. Return ONLY JSON: "
            '{"checks": [{"question": "...", "correct": true|false, "note": "..."}]}.\n\n'
            "Filled document text:\n" + filled_text[:8000] + "\n\n"
            "Question/answer pairs:\n" + "\n\n".join(qa_lines)
        )
        try:
            content = await llm.complete_json(prompt)
        except Exception:
            return {"status": "unavailable", "reason": "model error"}
        checks = _json_list_key(content, "checks")
        return {"status": "checked", "checks": checks}


doc_authoring_connector = DocAuthoringConnector()
