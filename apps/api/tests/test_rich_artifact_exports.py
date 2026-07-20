from __future__ import annotations

from io import BytesIO
import json
import os
import socket
import uuid
from zipfile import ZipFile
from contextlib import asynccontextmanager
import asyncio

import pytest

from core.artifact_rendering import build_preview, render_pdf_page
from core.artifact_rendering import safe_download_headers
from core.project_exports import _zip_bytes
from core.research_exports import build_research_docx, build_research_pdf


MAX_BYTES = 25 * 1024 * 1024
MAX_UNCOMPRESSED = 100 * 1024 * 1024


def _db_reachable() -> bool:
    host_port = os.environ.get(
        "DATABASE_URL", "postgresql+asyncpg://chronos:chronos@localhost:55432/chronos"
    ).rpartition("@")[2].partition("/")[0]
    host, _, port = host_port.rpartition(":")
    try:
        with socket.create_connection((host or "localhost", int(port or "55432")), timeout=1):
            return True
    except OSError:
        return False


_requires_db = pytest.mark.skipif(not _db_reachable(), reason="Postgres not reachable")


def preview(content: bytes, *, title: str, mime: str, kind: str = "file") -> dict:
    return build_preview(
        {"title": title, "mime_type": mime, "kind": kind},
        content,
        max_bytes=MAX_BYTES,
        max_uncompressed_bytes=MAX_UNCOMPRESSED,
        max_pdf_pages=50,
    )


def test_markup_and_executable_sources_are_inert() -> None:
    result = preview(
        b'<h1 onclick="steal()">Safe</h1><script>steal()</script>'
        b'<img src="https://tracker.example/x" onerror="steal()">',
        title="page.html",
        mime="text/html",
    )
    assert result["renderer"] == "markup"
    assert "Safe" in result["html"]
    assert "steal" not in result["html"]
    assert "tracker.example" not in result["html"]

    react = preview(
        b"export default function App(){ fetch('/secret'); return <main>Hello</main> }",
        title="App.tsx",
        mime="text/plain",
        kind="react",
    )
    assert react["renderer"] == "source"
    assert "without execution" in " ".join(react["limitations"])


def test_docx_pptx_and_xlsx_previews_extract_data_without_execution() -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    docx_out = BytesIO()
    document = Document()
    document.add_heading("Quarterly review", 1)
    document.add_paragraph("Revenue increased.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    document.save(docx_out)
    docx = preview(docx_out.getvalue(), title="review.docx", mime="application/octet-stream")
    assert docx["renderer"] == "document"
    assert docx["blocks"][0]["text"] == "Quarterly review"
    assert docx["tables"][0][0] == ["Metric", "Value"]

    pptx_out = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Launch plan"
    slide.placeholders[1].text = "Ship safely"
    presentation.save(pptx_out)
    pptx = preview(pptx_out.getvalue(), title="plan.pptx", mime="application/octet-stream")
    assert pptx["renderer"] == "presentation"
    assert "Launch plan" in pptx["slides"][0]["texts"]

    xlsx_out = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Forecast"
    sheet.append(["Revenue", 42])
    sheet["B2"] = "=1+1"
    workbook.save(xlsx_out)
    xlsx = preview(xlsx_out.getvalue(), title="forecast.xlsx", mime="application/octet-stream")
    assert xlsx["renderer"] == "workbook"
    assert xlsx["sheets"][0]["rows"][0][:2] == ["Revenue", "42"]
    assert "=1+1" not in json.dumps(xlsx)
    assert "never executed" in " ".join(xlsx["limitations"])


def test_notebook_and_zip_previews_never_execute_or_extract() -> None:
    notebook = {
        "nbformat": 4,
        "nbformat_minor": 5,
        "cells": [
            {
                "cell_type": "code",
                "source": ["print('safe')"],
                "outputs": [
                    {
                        "output_type": "display_data",
                        "data": {
                            "text/plain": ["safe output"],
                            "text/html": ["<script>steal()</script>"],
                        },
                    }
                ],
            }
        ],
    }
    ipynb = preview(
        json.dumps(notebook).encode(),
        title="analysis.ipynb",
        mime="application/x-ipynb+json",
    )
    assert ipynb["renderer"] == "notebook"
    assert ipynb["cells"][0]["outputs"] == ["safe output"]
    assert "<script>" not in json.dumps(ipynb)
    assert "never executed" in " ".join(ipynb["limitations"])

    archive_out = BytesIO()
    with ZipFile(archive_out, "w") as archive:
        archive.writestr("src/app.tsx", "export default 1")
    archive = preview(archive_out.getvalue(), title="project.zip", mime="application/zip")
    assert archive["renderer"] == "archive"
    assert archive["entries"][0]["path"] == "src/app.tsx"
    assert "never extracted" in " ".join(archive["limitations"])

    unsafe_out = BytesIO()
    with ZipFile(unsafe_out, "w") as unsafe:
        unsafe.writestr("../escape.txt", "no")
    rejected = preview(unsafe_out.getvalue(), title="unsafe.zip", mime="application/zip")
    assert rejected["status"] == "error"
    assert "unsafe path" in " ".join(rejected["limitations"])


def test_pdf_is_rasterized_and_research_exports_preserve_evidence() -> None:
    from docx import Document
    from pypdf import PdfReader
    from reportlab.pdfgen.canvas import Canvas

    pdf_source = BytesIO()
    canvas = Canvas(pdf_source)
    canvas.drawString(72, 720, "Safe PDF")
    canvas.showPage()
    canvas.save()
    source_bytes = pdf_source.getvalue()
    pdf_preview = preview(source_bytes, title="safe.pdf", mime="application/pdf")
    assert pdf_preview["renderer"] == "pdf"
    assert pdf_preview["page_count"] == 1
    assert render_pdf_page(source_bytes, 0, max_pages=50).startswith(b"\x89PNG")

    run = {
        "question": "What should Acme ship next?",
        "depth": "standard",
        "citation_policy": "required",
        "completed_at": "2026-07-13T12:00:00+00:00",
        "limitations": "Only public filings were available.",
    }
    report = "# Recommendation\n\nShip the reliable option [S1].\n\n## Sources\n\nOld duplicate."
    citations = [
        {
            "marker": "[S1]",
            "source_type": "web",
            "source_title": "Acme filing",
            "url": "https://example.com/filing",
            "snippet": "Reliability was the top customer concern.",
        }
    ]
    docx_bytes = build_research_docx(run, report, citations)
    docx_text = "\n".join(p.text for p in Document(BytesIO(docx_bytes)).paragraphs)
    assert "Ship the reliable option [S1]." in docx_text
    assert "Only public filings were available." in docx_text
    assert "Reliability was the top customer concern." in docx_text
    assert "Old duplicate." not in docx_text

    exported_pdf = build_research_pdf(run, report, citations)
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(exported_pdf)).pages)
    assert "Ship the reliable option [S1]." in pdf_text
    assert "Only public filings were available." in pdf_text
    assert "Reliability was the top customer concern." in pdf_text


def test_preview_size_limit_returns_honest_unsupported_state() -> None:
    result = build_preview(
        {"title": "huge.txt", "mime_type": "text/plain", "kind": "file"},
        b"x" * 20,
        max_bytes=10,
        max_uncompressed_bytes=100,
        max_pdf_pages=5,
    )
    assert result["status"] == "unsupported"
    assert result["renderer"] == "download"
    assert "download remains available" in result["limitations"][0]


def test_preview_response_and_pdf_raster_dimensions_are_bounded() -> None:
    from PIL import Image
    from reportlab.pdfgen.canvas import Canvas

    notebook = {
        "nbformat": 4,
        "cells": [
            {
                "cell_type": "code",
                "source": ["x" * 20_000],
                "outputs": [{"output_type": "stream", "text": ["y" * 20_000]}],
            }
            for _ in range(100)
        ],
    }
    bounded = preview(
        json.dumps(notebook).encode(),
        title="large.ipynb",
        mime="application/x-ipynb+json",
    )
    assert len(json.dumps(bounded)) < 1_100_000
    assert "1,000,000 characters" in " ".join(bounded["limitations"])

    source = BytesIO()
    canvas = Canvas(source, pagesize=(100_000, 100_000))
    canvas.drawString(72, 72, "bounded")
    canvas.showPage()
    canvas.save()
    png = render_pdf_page(source.getvalue(), 0, max_pages=2)
    with Image.open(BytesIO(png)) as image:
        assert image.width <= 4_096
        assert image.height <= 4_096
        assert image.width * image.height <= 16_000_000

    image_source = BytesIO()
    Image.new("RGB", (64, 32), "#224466").save(image_source, format="PNG")
    image_preview = preview(image_source.getvalue(), title="safe.png", mime="image/png")
    assert image_preview["renderer"] == "image"
    assert (image_preview["width"], image_preview["height"]) == (64, 32)

    invalid_image = preview(b"not really an image", title="unsafe.png", mime="image/png")
    assert invalid_image["status"] == "error"
    assert invalid_image["renderer"] == "download"


def test_project_bundle_manifest_and_download_filename_are_safe() -> None:
    import hashlib

    content = b"portable evidence"
    bundle = _zip_bytes(
        {"id": "project-1", "name": "Client launch"},
        [
            (
                {
                    "id": "artifact-1",
                    "title": "../../client\r\nreport.html",
                    "kind": "html",
                    "mime_type": "text/html",
                    "version": 3,
                    "created_at": "2026-07-16T00:00:00Z",
                    "updated_at": "2026-07-16T01:00:00Z",
                },
                content,
            )
        ],
    )
    with ZipFile(BytesIO(bundle)) as archive:
        names = archive.namelist()
        assert "manifest.json" in names
        artifact_name = next(name for name in names if name.startswith("artifacts/"))
        assert ".." not in artifact_name
        assert "\r" not in artifact_name and "\n" not in artifact_name
        manifest = json.loads(archive.read("manifest.json"))
        assert manifest["scope"] == "artifacts explicitly shared to this project"
        assert manifest["artifact_count"] == 1
        assert manifest["artifacts"][0]["sha256"] == hashlib.sha256(content).hexdigest()

    headers = safe_download_headers('bad"\r\nX-Evil: yes.html', active_markup=True)
    assert "\r" not in headers["Content-Disposition"]
    assert "\n" not in headers["Content-Disposition"]
    assert headers["Content-Disposition"].startswith("attachment;")
    assert headers["Content-Security-Policy"] == "default-src 'none'; sandbox"


@pytest.mark.asyncio
async def test_concurrent_research_exports_persist_exactly_one_artifact(monkeypatch) -> None:
    """The post-lock recheck prevents duplicates across simultaneous requests."""
    from core import research_exports

    report_id = "report-1"
    org_id = "tenant-1"
    persisted: dict[str, object] = {}
    save_calls = 0
    lock = asyncio.Lock()

    class FakeConnection:
        async def commit(self):
            return None

    connection = FakeConnection()

    async def fake_get_artifact(artifact_id: str):
        if artifact_id == report_id:
            return {
                "id": report_id,
                "organization_id": org_id,
                "is_deleted": False,
            }
        return persisted.get(artifact_id)

    async def fake_find(*_args, **_kwargs):
        return dict(persisted["export-1"]) if "export-1" in persisted else None

    @asynccontextmanager
    async def fake_lock(*_args, **_kwargs):
        async with lock:
            yield connection

    async def fake_save(_content: bytes, **kwargs):
        nonlocal save_calls
        save_calls += 1
        await asyncio.sleep(0)
        persisted["export-1"] = {
            "id": "export-1",
            "organization_id": kwargs["org_id"],
            "parent_artifact_id": kwargs["parent_artifact_id"],
            "kind": kwargs["kind"],
            "project_id": None,
            "is_deleted": False,
        }
        return "export-1"

    async def fake_content(_artifact_id: str):
        return b"# Result\n\nEvidence [S1]."

    async def fake_reflect(_name: str):
        return object()

    async def fake_set_project(_artifacts, artifact_id: str, **_kwargs):
        return dict(persisted[artifact_id])

    monkeypatch.setattr(research_exports, "get_artifact", fake_get_artifact)
    monkeypatch.setattr(research_exports, "_find_existing_export", fake_find)
    monkeypatch.setattr(research_exports, "_export_lock", fake_lock)
    monkeypatch.setattr(research_exports, "save_artifact", fake_save)
    monkeypatch.setattr(research_exports, "read_artifact_content", fake_content)
    monkeypatch.setattr(research_exports, "reflect_table", fake_reflect)
    monkeypatch.setattr(research_exports, "_set_project_and_fetch", fake_set_project)
    monkeypatch.setattr(research_exports, "build_research_pdf", lambda *_args: b"%PDF-safe")

    run = {
        "status": "complete",
        "report_artifact_id": report_id,
        "question": "Concurrent export?",
    }
    results = await asyncio.gather(
        *(
            research_exports.create_research_export(
                run, [], "pdf", org_id=org_id, created_by="member:1"
            )
            for _ in range(8)
        )
    )

    assert save_calls == 1
    assert {result[0]["id"] for result in results} == {"export-1"}
    assert sum(1 for _, reused in results if not reused) == 1


@_requires_db
@pytest.mark.asyncio
async def test_postgres_export_lock_is_exactly_once_across_real_connections() -> None:
    from sqlalchemy import func, select

    from core.artifacts import save_artifact
    from core.db import engine, reflect_table
    from core.research_exports import create_research_export

    org_id = f"artifact-export-{uuid.uuid4().hex[:10]}"
    report_id = await save_artifact(
        "# Production evidence\n\nThe result is reproducible [S1].",
        kind="markdown",
        title="Production evidence.md",
        mime_type="text/markdown",
        org_id=org_id,
        created_by="member:integration",
    )
    run = {
        "status": "complete",
        "report_artifact_id": report_id,
        "question": "Is the export idempotent?",
        "citation_policy": "required",
        "limitations": "Integration-test scope only.",
    }
    citations = [
        {
            "marker": "[S1]",
            "source_type": "test",
            "source_title": "Deterministic fixture",
            "snippet": "The result is reproducible.",
        }
    ]
    results = await asyncio.gather(
        *(
            create_research_export(
                run,
                citations,
                "pdf",
                org_id=org_id,
                created_by="member:integration",
            )
            for _ in range(6)
        )
    )
    ids = {str(artifact["id"]) for artifact, _reused in results}
    assert len(ids) == 1
    assert sum(1 for _artifact, reused in results if not reused) == 1

    artifacts = await reflect_table("artifacts")
    async with engine.begin() as conn:
        count = await conn.scalar(
            select(func.count())
            .select_from(artifacts)
            .where(
                artifacts.c.organization_id == org_id,
                artifacts.c.parent_artifact_id == report_id,
                artifacts.c.kind == "research_report_pdf",
                artifacts.c.is_deleted == False,  # noqa: E712
            )
        )
    assert count == 1


@_requires_db
@pytest.mark.asyncio
async def test_project_bundle_contains_only_explicit_project_artifacts() -> None:
    from sqlalchemy import insert

    from core.artifacts import save_artifact, set_artifact_project
    from core.db import engine, reflect_table
    from core.project_exports import build_project_bundle

    org_id = f"project-export-{uuid.uuid4().hex[:10]}"
    projects = await reflect_table("projects")
    async with engine.begin() as conn:
        project_id = str(
            (
                await conn.execute(
                    insert(projects)
                    .values(organization_id=org_id, name="Explicit scope")
                    .returning(projects.c.id)
                )
            ).scalar_one()
        )
    shared_id = await save_artifact(
        "shared",
        kind="markdown",
        title="shared.md",
        org_id=org_id,
    )
    private_id = await save_artifact(
        "private task output",
        kind="markdown",
        title="private.md",
        task_id=str(uuid.uuid4()),
        org_id=org_id,
    )
    prior_bundle_id = await save_artifact(
        b"old bundle",
        kind="project_bundle",
        title="old.zip",
        mime_type="application/zip",
        org_id=org_id,
    )
    await set_artifact_project(shared_id, project_id=project_id, org_id=org_id)
    await set_artifact_project(prior_bundle_id, project_id=project_id, org_id=org_id)

    bundle, summary = await build_project_bundle(project_id, org_id)
    with ZipFile(BytesIO(bundle)) as archive:
        manifest = json.loads(archive.read("manifest.json"))
    exported_ids = {item["id"] for item in manifest["artifacts"]}
    assert exported_ids == {shared_id}
    assert private_id not in exported_ids
    assert prior_bundle_id not in exported_ids
    assert summary["scope"] == "explicit_project_artifacts"
