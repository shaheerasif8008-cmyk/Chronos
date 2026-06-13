"""Acceptance proof for the document authoring connector.

Covers doc.create (pdf/docx/markdown), doc.create_slides (pptx), doc.fill_pdf
(overlay on the original page, by explicit coordinates and by anchor text), and
doc.render_chart (matplotlib → image artifact).

Artifact storage is stubbed via monkeypatch so these tests need no database: the
connector is a pure renderer and we assert on the bytes it produces (re-parsing
them with the same libraries that read them in production).
"""
from __future__ import annotations

import io

import pytest

import connectors.doc_authoring as da


# ---------------------------------------------------------------------------
# Harness: capture saved artifacts and stub artifact resolution
# ---------------------------------------------------------------------------


class _Captured:
    def __init__(self):
        self.saved: list[dict] = []
        # artifact_id -> (bytes, org_id)
        self.store: dict[str, tuple[bytes, str]] = {}

    def add_source(self, artifact_id: str, raw: bytes, org_id: str = "default") -> None:
        self.store[artifact_id] = (raw, org_id)


@pytest.fixture
def cap(monkeypatch):
    c = _Captured()

    async def fake_save(content, *, kind, title=None, task_id=None, org_id="default",
                        mime_type=None, parent_artifact_id=None, created_by=None, **kw):
        raw = content if isinstance(content, bytes) else content.encode()
        aid = f"artifact-{len(c.saved)}"
        c.saved.append({
            "artifact_id": aid, "bytes": raw, "kind": kind, "title": title,
            "mime_type": mime_type, "org_id": org_id, "parent": parent_artifact_id,
        })
        return aid

    async def fake_get(artifact_id):
        if artifact_id not in c.store:
            return None
        _, org = c.store[artifact_id]
        return {"organization_id": org, "is_deleted": False}

    async def fake_read(artifact_id):
        if artifact_id not in c.store:
            return None
        return c.store[artifact_id][0]

    monkeypatch.setattr(da, "save_artifact", fake_save)
    monkeypatch.setattr(da, "get_artifact", fake_get)
    monkeypatch.setattr(da, "read_artifact_content", fake_read)
    return c


def _source_pdf(text: str, *, anchor: str = "Answer:") -> bytes:
    """Build a one-page PDF containing a line of text (acts as a worksheet)."""
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    cv = canvas.Canvas(buf, pagesize=(612, 792))
    cv.setFont("Helvetica", 12)
    cv.drawString(72, 700, text)
    cv.save()
    return buf.getvalue()


def _pdf_text(raw: bytes) -> str:
    from pypdf import PdfReader

    return "\n".join((p.extract_text() or "") for p in PdfReader(io.BytesIO(raw)).pages)


# ---------------------------------------------------------------------------
# doc.create
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_create_pdf(cap):
    result = await da.doc_authoring_connector.execute("doc.create", {
        "__org_id": "org1", "format": "pdf", "title": "Lab Report",
        "blocks": [
            {"type": "heading", "text": "Results", "level": 1},
            {"type": "paragraph", "text": "The reaction was exothermic."},
            {"type": "bullet", "text": "Observation one"},
        ],
    })
    assert result.data["status"] == "success"
    saved = cap.saved[-1]
    assert saved["mime_type"] == "application/pdf"
    assert saved["bytes"][:4] == b"%PDF"
    text = _pdf_text(saved["bytes"])
    assert "Results" in text and "exothermic" in text


@pytest.mark.asyncio
async def test_create_markdown(cap):
    result = await da.doc_authoring_connector.execute("doc.create", {
        "__org_id": "org1", "format": "markdown", "title": "Notes",
        "blocks": [{"type": "bullet", "text": "point A"}],
    })
    assert result.data["status"] == "success"
    body = cap.saved[-1]["bytes"].decode()
    assert "# Notes" in body and "- point A" in body


@pytest.mark.asyncio
async def test_create_docx(cap):
    result = await da.doc_authoring_connector.execute("doc.create", {
        "__org_id": "org1", "format": "docx", "title": "Memo",
        "blocks": [{"type": "paragraph", "text": "hello world"}],
    })
    assert result.data["status"] == "success"
    from docx import Document
    doc = Document(io.BytesIO(cap.saved[-1]["bytes"]))
    assert any("hello world" in p.text for p in doc.paragraphs)


@pytest.mark.asyncio
async def test_create_rejects_empty_blocks(cap):
    result = await da.doc_authoring_connector.execute("doc.create", {
        "__org_id": "org1", "format": "pdf", "title": "x", "blocks": [],
    })
    assert result.data["status"] == "error"
    assert not cap.saved


# ---------------------------------------------------------------------------
# doc.create_slides
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_create_slides(cap):
    result = await da.doc_authoring_connector.execute("doc.create_slides", {
        "__org_id": "org1", "title": "Deck",
        "slides": [
            {"title": "Intro", "bullets": ["first", "second"], "notes": "say hi"},
            {"title": "End"},
        ],
    })
    assert result.data["status"] == "success"
    assert result.data["slides"] == 2
    from pptx import Presentation
    prs = Presentation(io.BytesIO(cap.saved[-1]["bytes"]))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "Intro" in titles and "End" in titles


# ---------------------------------------------------------------------------
# doc.fill_pdf — overlay on the original page
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_fill_pdf_by_anchor(cap):
    src = _source_pdf("Q1: What is H2O? Answer:")
    cap.add_source("src-1", src, org_id="org1")

    result = await da.doc_authoring_connector.execute("doc.fill_pdf", {
        "__org_id": "org1", "artifact_id": "src-1",
        "items": [
            {"page": 1, "type": "text", "anchor_text": "Answer:",
             "dx": 90, "dy": 0, "text": "Water", "color": "#0000ff"},
        ],
    })
    assert result.data["status"] == "success"
    assert result.data["items_placed"] == 1
    out = cap.saved[-1]
    # Output preserves the original content AND carries the overlay.
    text = _pdf_text(out["bytes"])
    assert "What is H2O" in text  # original preserved
    assert "Water" in text         # overlay applied
    assert out["parent"] == "src-1"


@pytest.mark.asyncio
async def test_fill_pdf_by_explicit_coords(cap):
    src = _source_pdf("Name: ____")
    cap.add_source("src-2", src, org_id="org1")

    result = await da.doc_authoring_connector.execute("doc.fill_pdf", {
        "__org_id": "org1", "artifact_id": "src-2",
        "items": [{"page": 1, "type": "text", "x": 120, "y": 92, "text": "Ada Lovelace"}],
    })
    assert result.data["status"] == "success"
    assert "Ada Lovelace" in _pdf_text(cap.saved[-1]["bytes"])


@pytest.mark.asyncio
async def test_fill_pdf_skips_out_of_range_page(cap):
    src = _source_pdf("single page")
    cap.add_source("src-3", src, org_id="org1")

    result = await da.doc_authoring_connector.execute("doc.fill_pdf", {
        "__org_id": "org1", "artifact_id": "src-3",
        "items": [
            {"page": 1, "type": "text", "x": 72, "y": 100, "text": "ok"},
            {"page": 9, "type": "text", "x": 72, "y": 100, "text": "nope"},
        ],
    })
    assert result.data["status"] == "success"
    assert result.data["items_placed"] == 1
    assert result.data["items_skipped"]  # the page-9 item is reported, not silently dropped


@pytest.mark.asyncio
async def test_fill_pdf_unknown_anchor_places_nothing(cap):
    src = _source_pdf("nothing to match here")
    cap.add_source("src-4", src, org_id="org1")

    result = await da.doc_authoring_connector.execute("doc.fill_pdf", {
        "__org_id": "org1", "artifact_id": "src-4",
        "items": [{"page": 1, "type": "text", "anchor_text": "Zzz", "text": "x"}],
    })
    assert result.data["status"] == "success"
    assert result.data["items_placed"] == 0


@pytest.mark.asyncio
async def test_fill_pdf_cross_org_denied(cap):
    src = _source_pdf("private")
    cap.add_source("src-5", src, org_id="orgA")

    result = await da.doc_authoring_connector.execute("doc.fill_pdf", {
        "__org_id": "orgB", "artifact_id": "src-5",
        "items": [{"page": 1, "type": "text", "x": 72, "y": 100, "text": "x"}],
    })
    assert result.data["status"] == "error"
    assert not cap.saved


@pytest.mark.asyncio
async def test_fill_pdf_non_pdf_source(cap):
    cap.add_source("src-6", b"this is not a pdf", org_id="org1")
    result = await da.doc_authoring_connector.execute("doc.fill_pdf", {
        "__org_id": "org1", "artifact_id": "src-6",
        "items": [{"page": 1, "type": "text", "x": 1, "y": 1, "text": "x"}],
    })
    assert result.data["status"] == "error"


# ---------------------------------------------------------------------------
# doc.render_chart
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_render_chart_line(cap):
    result = await da.doc_authoring_connector.execute("doc.render_chart", {
        "__org_id": "org1", "chart_type": "line", "title": "Trend",
        "series": [{"name": "a", "values": [1, 2, 3]}],
        "labels": ["x", "y", "z"],
    })
    assert result.data["status"] == "success"
    saved = cap.saved[-1]
    assert saved["kind"] == "image"
    assert saved["bytes"][:4] == b"\x89PNG"


@pytest.mark.asyncio
async def test_render_chart_pie(cap):
    result = await da.doc_authoring_connector.execute("doc.render_chart", {
        "__org_id": "org1", "chart_type": "pie",
        "values": [10, 20, 30], "labels": ["a", "b", "c"],
    })
    assert result.data["status"] == "success"
    assert cap.saved[-1]["bytes"][:4] == b"\x89PNG"


# ---------------------------------------------------------------------------
# doc.detect_fields — vision auto-coordinate detection
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_detect_fields_converts_to_points(cap, monkeypatch):
    src = _source_pdf("Name: ____   Age: ____")
    cap.add_source("form-1", src, org_id="org1")

    async def fake_vision_json(image_bytes, mime, instruction):
        # Two normalized boxes (top-left origin), as a vision model would return.
        return (
            '{"fields": [{"label": "Name", "field_type": "text", "bbox": [0.2, 0.1, 0.5, 0.15]},'
            ' {"label": "Age", "field_type": "text", "bbox": [0.6, 0.1, 0.8, 0.15]}]}'
        )

    import core.llm as llm_mod
    monkeypatch.setattr(llm_mod, "vision_json", fake_vision_json)

    result = await da.doc_authoring_connector.execute("doc.detect_fields", {
        "__org_id": "org1", "artifact_id": "form-1",
    })
    assert result.data["status"] == "success"
    fields = result.data["fields"]
    assert len(fields) == 2
    f = fields[0]
    # Page is 612pt wide, 792 tall; x0=0.2 → ~122.4pt, y0=0.1 → ~79.2pt (top-left origin).
    assert f["page"] == 1
    assert abs(f["bbox_points"]["x"] - 0.2 * 612) < 1.0
    assert abs(f["bbox_points"]["y"] - 0.1 * 792) < 1.0
    # fill_hint is ready to drop into doc.fill_pdf.
    assert "x" in f["fill_hint"] and "y" in f["fill_hint"] and "size" in f["fill_hint"]


@pytest.mark.asyncio
async def test_detect_fields_no_vision_model_degrades(cap, monkeypatch):
    src = _source_pdf("Q: ____")
    cap.add_source("form-2", src, org_id="org1")

    async def empty_vision(image_bytes, mime, instruction):
        return ""  # no vision model configured

    import core.llm as llm_mod
    monkeypatch.setattr(llm_mod, "vision_json", empty_vision)

    result = await da.doc_authoring_connector.execute("doc.detect_fields", {
        "__org_id": "org1", "artifact_id": "form-2",
    })
    assert result.data["status"] == "unavailable"
    assert result.data["fields"] == []
    assert "fallback_reason" in result.data


@pytest.mark.asyncio
async def test_detect_fields_cross_org_denied(cap, monkeypatch):
    src = _source_pdf("x")
    cap.add_source("form-3", src, org_id="orgA")
    result = await da.doc_authoring_connector.execute("doc.detect_fields", {
        "__org_id": "orgB", "artifact_id": "form-3",
    })
    assert result.data["status"] == "error"


# ---------------------------------------------------------------------------
# doc.verify_fill — re-read the filled PDF and confirm the work
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_verify_fill_detects_present_and_missing(cap):
    # A "filled" PDF whose text contains one of the two expected answers.
    filled = _source_pdf("Answer: Water is H2O")
    cap.add_source("filled-1", filled, org_id="org1")

    result = await da.doc_authoring_connector.execute("doc.verify_fill", {
        "__org_id": "org1", "artifact_id": "filled-1",
        "expected": [
            {"label": "Q1", "answer": "Water is H2O"},
            {"label": "Q2", "answer": "Sodium Chloride"},
        ],
    })
    assert result.data["status"] == "success"
    assert result.data["all_present"] is False
    assert result.data["missing"] == ["Sodium Chloride"]
    by_label = {it["label"]: it["present"] for it in result.data["items"]}
    assert by_label["Q1"] is True and by_label["Q2"] is False


@pytest.mark.asyncio
async def test_verify_fill_all_present(cap):
    filled = _source_pdf("Paris is the capital")
    cap.add_source("filled-2", filled, org_id="org1")
    result = await da.doc_authoring_connector.execute("doc.verify_fill", {
        "__org_id": "org1", "artifact_id": "filled-2",
        "expected": [{"answer": "Paris"}],
    })
    assert result.data["all_present"] is True
    assert result.data["missing"] == []


@pytest.mark.asyncio
async def test_verify_fill_correctness_recheck(cap, monkeypatch):
    filled = _source_pdf("Answer: 4")
    cap.add_source("filled-3", filled, org_id="org1")

    async def fake_complete_json(prompt, **kw):
        return '{"checks": [{"question": "2+2?", "correct": true, "note": "ok"}]}'

    import core.llm as llm_mod
    monkeypatch.setattr(llm_mod, "complete_json", fake_complete_json)

    result = await da.doc_authoring_connector.execute("doc.verify_fill", {
        "__org_id": "org1", "artifact_id": "filled-3",
        "expected": [{"answer": "4", "question": "2+2?"}],
        "recheck_correctness": True,
    })
    assert result.data["correctness"]["status"] == "checked"
    assert result.data["correctness"]["checks"][0]["correct"] is True


@pytest.mark.asyncio
async def test_verify_fill_requires_expected(cap):
    cap.add_source("filled-4", _source_pdf("x"), org_id="org1")
    result = await da.doc_authoring_connector.execute("doc.verify_fill", {
        "__org_id": "org1", "artifact_id": "filled-4", "expected": [],
    })
    assert result.data["status"] == "error"


@pytest.mark.asyncio
async def test_unknown_tool_raises(cap):
    with pytest.raises(ValueError):
        await da.doc_authoring_connector.execute("doc.bogus", {"__org_id": "org1"})
