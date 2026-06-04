import base64
import io
import os
from unittest.mock import AsyncMock, patch

import pytest

from core import llm
from core.models import AgentContext
from parsing.engine import PREVIEW_CHAR_LIMIT, ParsedDocument, parse_document

# ---------------------------------------------------------------------------
# DB connectivity guard — used by the storage round-trip test only.
# ---------------------------------------------------------------------------
def _db_reachable() -> bool:
    import socket

    host, _, port_str = os.environ.get(
        "DATABASE_URL", "postgresql+asyncpg://chronos:chronos@localhost:5432/chronos"
    ).rpartition("@")[-1].partition("/")[0].rpartition(":")
    try:
        with socket.create_connection((host or "localhost", int(port_str or 5432)), timeout=1):
            return True
    except OSError:
        return False


_requires_db = pytest.mark.skipif(not _db_reachable(), reason="Postgres not reachable")


@pytest.fixture(autouse=True)
def dispose_db_engine():
    """Dispose the SQLAlchemy engine pool before each test in this module.

    pytest-asyncio gives each test its own event loop (function scope). asyncpg
    connection pools are bound to the loop that created them, so a pool created
    in test N's loop is unusable in test N+1's loop. Disposing before each test
    forces a fresh pool in the current event loop, preventing 'Future attached
    to a different loop' / 'Event loop is closed' failures when DB tests follow
    non-DB tests.
    """
    import core.db as _db
    _db.engine.sync_engine.pool.dispose()
    yield


@pytest.mark.asyncio
async def test_vision_ocr_returns_empty_when_no_model_configured():
    with patch.object(llm.settings, "vision_model", ""):
        out = await llm.vision_ocr(b"\x89PNG\r\n", "image/png")
    assert out == ""


@pytest.mark.asyncio
async def test_vision_ocr_sends_data_url_and_returns_text():
    fake = {"choices": [{"message": {"content": "INVOICE 42"}}]}
    with patch.object(llm.settings, "vision_model", "openrouter/openai/gpt-4o-mini"), \
         patch.object(llm, "litellm") as mock_litellm:
        mock_litellm.acompletion = AsyncMock(return_value=fake)
        out = await llm.vision_ocr(b"rawbytes", "image/png")
    assert out == "INVOICE 42"
    sent = mock_litellm.acompletion.call_args.kwargs
    content = sent["messages"][0]["content"]
    # multimodal: a text part + an image_url data URL part
    image_part = next(p for p in content if p["type"] == "image_url")
    assert image_part["image_url"]["url"].startswith("data:image/png;base64,")
    assert base64.b64encode(b"rawbytes").decode() in image_part["image_url"]["url"]


@pytest.mark.asyncio
async def test_parse_plain_text():
    doc = await parse_document(b"hello world", "text/plain", "note.txt")
    assert isinstance(doc, ParsedDocument)
    assert doc.full_text == "hello world"
    assert doc.parser_used == "text"
    assert doc.truncated is False
    assert doc.note is None


@pytest.mark.asyncio
async def test_parse_csv():
    doc = await parse_document(b"a,b\n1,2\n", "text/csv", "data.csv")
    assert "a,b" in doc.full_text
    assert doc.parser_used == "text"


@pytest.mark.asyncio
async def test_parse_unknown_type_is_unparseable_not_crash():
    doc = await parse_document(b"\x00\x01\x02", "application/x-thing", "blob.bin")
    assert doc.parser_used == "none"
    assert doc.note == "not parseable yet"
    assert doc.full_text == ""


@pytest.mark.asyncio
async def test_preview_truncates_long_text():
    big = "x" * (PREVIEW_CHAR_LIMIT + 500)
    doc = await parse_document(big.encode(), "text/plain", "big.txt")
    assert doc.truncated is True
    assert len(doc.preview) == PREVIEW_CHAR_LIMIT
    assert doc.full_text == big


def _one_page_pdf_with_text(text: str) -> bytes:
    from pypdf import PdfWriter
    # pypdf can't author text-bearing pages easily; use reportlab-free minimal route
    # via a blank page, then monkeypatch extraction in the test instead.
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


@pytest.mark.asyncio
async def test_parse_pdf_extracts_text_layer():
    pdf = _one_page_pdf_with_text("ignored")
    with patch("parsing.engine._pdf_page_texts", return_value=["Quarterly report", "page two"]):
        doc = await parse_document(pdf, "application/pdf", "report.pdf")
    assert "Quarterly report" in doc.full_text
    assert "page two" in doc.full_text
    assert doc.page_count == 2
    assert doc.parser_used == "pdf"


@pytest.mark.asyncio
async def test_parse_pdf_falls_back_to_ocr_on_empty_pages():
    pdf = _one_page_pdf_with_text("ignored")
    with patch("parsing.engine._pdf_page_texts", return_value=["", ""]), \
         patch("parsing.engine._pdf_page_ocr", new=AsyncMock(return_value="scanned text")):
        doc = await parse_document(pdf, "application/pdf", "scan.pdf")
    assert "scanned text" in doc.full_text
    assert doc.parser_used == "pdf+ocr"


@pytest.mark.asyncio
async def test_parse_image_uses_vision_ocr():
    with patch("parsing.engine.vision_ocr", new=AsyncMock(return_value="receipt total $9")):
        doc = await parse_document(b"\x89PNG", "image/png", "receipt.png")
    assert doc.full_text == "receipt total $9"
    assert doc.parser_used == "image-ocr"


def _make_docx(text: str) -> bytes:
    from docx import Document
    d = Document()
    d.add_paragraph(text)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _make_xlsx(value: str) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    wb.active["A1"] = value
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.mark.asyncio
async def test_parse_docx():
    doc = await parse_document(_make_docx("Contract clause 7"), "", "agreement.docx")
    assert "Contract clause 7" in doc.full_text
    assert doc.parser_used == "docx"


@pytest.mark.asyncio
async def test_parse_xlsx():
    doc = await parse_document(_make_xlsx("Revenue 2026"), "", "model.xlsx")
    assert "Revenue 2026" in doc.full_text
    assert doc.parser_used == "xlsx"


@pytest.mark.asyncio
async def test_parse_pptx():
    doc = await parse_document(_make_pptx("Roadmap Q3"), "", "deck.pptx")
    assert "Roadmap Q3" in doc.full_text
    assert doc.parser_used == "pptx"


@pytest.mark.asyncio
async def test_parse_docx_with_table():
    from docx import Document as DocxDocument
    d = DocxDocument()
    t = d.add_table(rows=2, cols=2)
    t.rows[0].cells[0].text = "Name"
    t.rows[0].cells[1].text = "Score"
    t.rows[1].cells[0].text = "Alice"
    t.rows[1].cells[1].text = "95"
    buf = io.BytesIO()
    d.save(buf)
    doc = await parse_document(buf.getvalue(), "", "data.docx")
    assert "Name | Score" in doc.full_text
    assert "Alice | 95" in doc.full_text


@pytest.mark.asyncio
async def test_parse_xlsx_multi_sheet():
    from openpyxl import Workbook
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Revenue"
    ws1["A1"] = "Q1"
    ws2 = wb.create_sheet("Costs")
    ws2["A1"] = "Q2"
    buf = io.BytesIO()
    wb.save(buf)
    doc = await parse_document(buf.getvalue(), "", "model.xlsx")
    assert "# Sheet: Revenue" in doc.full_text
    assert "# Sheet: Costs" in doc.full_text
    assert doc.page_count == 2


@pytest.mark.asyncio
async def test_parse_pptx_multi_slide():
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Intro"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Body"
    buf = io.BytesIO()
    prs.save(buf)
    doc = await parse_document(buf.getvalue(), "", "deck.pptx")
    assert "# Slide 1" in doc.full_text
    assert "# Slide 2" in doc.full_text
    assert doc.page_count == 2


@pytest.mark.asyncio
async def test_parse_docx_corrupt_bytes():
    doc = await parse_document(b"not a zip", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "bad.docx")
    assert doc.parser_used == "none"
    assert doc.note is not None
    assert "docx" in doc.note


@_requires_db
@pytest.mark.asyncio
async def test_save_artifact_records_parent_and_status():
    from core import artifacts

    parent = await artifacts.save_artifact(
        b"raw pdf bytes", kind="attachment", title="report.pdf",
        mime_type="application/pdf", parse_status="pending",
    )
    child = await artifacts.save_artifact(
        "extracted text", kind="parsed_text", title="report.pdf (text)",
        parent_artifact_id=parent, parse_status="parsed",
    )
    meta = await artifacts.get_artifact(child)
    assert str(meta["parent_artifact_id"]) == parent
    assert meta["parse_status"] == "parsed"


def test_doc_tools_registered_and_broker_named():
    from runtime.tool_registry import ALL_TOOLS, SUBAGENT_TOOLS, to_broker_name, tool_name

    names = {tool_name(s) for s in ALL_TOOLS}
    assert "doc__parse" in names
    assert "doc__read" in names
    # available to sub-agents too
    sub_names = {tool_name(s) for s in SUBAGENT_TOOLS}
    assert "doc__parse" in sub_names and "doc__read" in sub_names
    # convert cleanly to broker dot-notation
    assert to_broker_name("doc__parse") == "doc.parse"
    assert to_broker_name("doc__read") == "doc.read"


@pytest.mark.asyncio
async def test_doc_connector_parse_by_path(tmp_path, monkeypatch):
    from parsing import tool as doctool

    # Point the workspace root at a temp dir and drop a file in it.
    monkeypatch.setattr("connectors.filesystem.WORKSPACE_ROOT", tmp_path)
    work = tmp_path / "default" / "t1"
    work.mkdir(parents=True)
    (work / "note.txt").write_text("workspace doc body")

    args = {"path": "note.txt", "__org_id": "default", "__task_id": "t1"}
    result = await doctool.doc_connector.execute("doc.parse", args)
    assert "workspace doc body" in result.data["preview"]
    assert result.data["parser_used"] == "text"


@pytest.mark.asyncio
async def test_doc_connector_read_pages_artifact(monkeypatch):
    from parsing import tool as doctool

    async def fake_read(artifact_id):
        return b"A" * 100

    async def fake_meta(artifact_id):
        return {"mime_type": "text/plain", "title": "big.txt", "organization_id": "default", "kind": "attachment"}

    monkeypatch.setattr(doctool, "read_artifact_content", fake_read)
    monkeypatch.setattr(doctool, "get_artifact", fake_meta)

    args = {"artifact_id": "x", "char_offset": 10, "max_chars": 5, "__org_id": "default", "__task_id": "t1"}
    result = await doctool.doc_connector.execute("doc.read", args)
    assert result.data["content"] == "AAAAA"
    assert result.data["char_offset"] == 10


@pytest.mark.asyncio
async def test_upload_attachment_stores_and_returns_id(monkeypatch):
    from routers import attachments as att_router
    from io import BytesIO
    from starlette.datastructures import UploadFile as StarletteUploadFile, Headers

    saved: dict = {}

    async def fake_save(raw, **kw):
        saved.update(kw)
        saved["raw"] = raw
        return "att-123"

    async def fake_log(*a, **k):
        return None

    async def fake_check(*a, **k):
        return True

    async def fake_require_conv(*a, **k):
        pass  # bypass real DB lookup for unit test

    monkeypatch.setattr(att_router, "save_artifact", fake_save)
    monkeypatch.setattr(att_router.audit, "log", fake_log)
    monkeypatch.setattr(att_router.permissions, "check", fake_check)
    monkeypatch.setattr(att_router, "_require_conversation_member", fake_require_conv)

    # Adjust Member constructor kwargs to match core/models.py exactly.
    from core.models import Member
    member = Member(id="m1", organization_id="default", email="a@b.c", role="user", name="A")
    upload = StarletteUploadFile(
        filename="report.pdf",
        file=BytesIO(b"%PDF-1.4 data"),
        headers=Headers({"content-type": "application/pdf"}),
    )
    out = await att_router.upload_attachment(
        file=upload, conversation_id="c1", project_id=None, task_id=None,
        research_run_id=None, member=member,
    )
    assert out["attachment_id"] == "att-123"
    assert out["filename"] == "report.pdf"
    assert saved["kind"] == "attachment"
    assert saved["parse_status"] == "pending"


@pytest.mark.asyncio
async def test_doc_parse_routes_through_broker(monkeypatch):
    from core import tool_broker as tb
    from core.models import ToolResult

    audited: list[str] = []

    async def fake_log(event_type, actor, action, **kw):
        audited.append(event_type)

    async def fake_check(*a, **k):
        return True

    async def fake_tool_policy(*a, **k):
        return {}

    monkeypatch.setattr(tb.audit, "log", fake_log)
    monkeypatch.setattr(tb.permissions, "check", fake_check)
    monkeypatch.setattr(tb, "tool_policy", fake_tool_policy)
    # connector_tier is imported into tool_broker's namespace; patch it there.
    monkeypatch.setattr(tb, "connector_tier", AsyncMock(return_value="live"))

    async def fake_doc_exec(tool, args):
        return ToolResult(data={"preview": "hi"}, summary="ok")

    monkeypatch.setattr("parsing.tool.doc_connector.execute", fake_doc_exec)

    agent = AgentContext(id="a1", org_id="default", task_id="t1", member_id="m1")
    result = await tb.execute(agent, "doc.parse", {"artifact_id": "x"})
    assert result.summary == "ok"
    assert "tool_call" in audited and "tool_result" in audited


@pytest.mark.asyncio
async def test_load_history_injects_attachments_block():
    from runtime import agent_loop

    task = {
        "id": "t1",
        "goal": "Summarize the attached report",
        "agent_state": {
            "agent_history": [],
            "attachments": [
                {"filename": "report.pdf", "preview": "Q3 revenue up 12%", "truncated": True, "parsed_artifact_id": "p1"},
            ],
        },
    }
    history = await agent_loop._load_history(task, tools=[])
    blocks = [m["content"] for m in history if m["role"] == "user"]
    assert any("Q3 revenue up 12%" in b and "report.pdf" in b for b in blocks)
    assert any("doc__read" in b for b in blocks)  # truncation hint present
    assert history[-1]["content"] == "Summarize the attached report"  # goal is last


@pytest.mark.asyncio
async def test_parse_attachments_sets_status_and_returns_preview(monkeypatch):
    from routers import chat as chat_router

    async def fake_get(artifact_id):
        return {"mime_type": "text/plain", "title": "note.txt", "organization_id": "default", "parse_status": "pending"}

    async def fake_read(artifact_id):
        return b"important content"

    async def fake_save(*a, **kw):
        return "parsed-001"

    async def fake_set_status(artifact_id, status):
        pass

    async def fake_audit_log(*a, **kw):
        pass

    monkeypatch.setattr(chat_router, "_get_artifact", fake_get)
    monkeypatch.setattr(chat_router, "_read_artifact_content", fake_read)
    monkeypatch.setattr(chat_router, "_save_artifact", fake_save)
    monkeypatch.setattr(chat_router, "_set_parse_status", fake_set_status)
    monkeypatch.setattr(chat_router.audit, "log", fake_audit_log)

    result = await chat_router._parse_attachments(["att-1"], "conv-1", "default")
    assert len(result) == 1
    assert result[0]["preview"] == "important content"
    assert result[0]["filename"] == "note.txt"
    assert result[0]["parsed_artifact_id"] == "parsed-001"


@pytest.mark.asyncio
async def test_attachment_not_parseable_across_orgs(monkeypatch):
    """An attachment owned by org A must not be parsed when requested as org B."""
    from routers import chat

    async def fake_get_artifact(att_id):
        return {"id": att_id, "organization_id": "orgA", "mime_type": "text/plain", "title": "secret.txt"}

    monkeypatch.setattr(chat, "_get_artifact", fake_get_artifact)

    out = await chat._parse_attachments(["att-A"], conversation_id="c1", org_id="orgB")
    assert out == []  # org mismatch → skipped, no parsed_text created


@pytest.mark.asyncio
async def test_doc_connector_rejects_cross_org_artifact(monkeypatch):
    """DocConnector must reject artifact reads for a different org."""
    from parsing import tool as doctool

    async def fake_meta(artifact_id):
        return {"organization_id": "orgA", "mime_type": "text/plain", "title": "secret.txt", "kind": "attachment"}

    monkeypatch.setattr(doctool, "get_artifact", fake_meta)

    with pytest.raises(PermissionError):
        await doctool.doc_connector.execute("doc.parse", {"artifact_id": "x", "__org_id": "orgB"})

    with pytest.raises(PermissionError):
        await doctool.doc_connector.execute("doc.read", {"artifact_id": "x", "__org_id": "orgB"})
