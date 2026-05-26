"""Pure document parsing engine.

bytes + mime + filename → ParsedDocument. No DB, no object storage, no network
except the litellm OCR call. Both the chat-upload surface and the agent doc tools
wrap this module.
"""
from __future__ import annotations

import io
from dataclasses import dataclass

from core.llm import vision_ocr

#: Preview is the first ~6K tokens, approximated at 4 chars/token.
PREVIEW_CHAR_LIMIT = 24_000

#: Note set when the file type is not supported (vs a recognized type that errored).
#: Used by callers to distinguish "unparseable" from "failed".
UNPARSEABLE_NOTE = "not parseable yet"

_TEXT_MIMES = {"text/plain", "text/csv", "text/markdown", "application/csv"}
_TEXT_EXTS = {".txt", ".csv", ".md", ".markdown", ".log", ".tsv"}
_IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg", "image/webp"}


@dataclass
class ParsedDocument:
    full_text: str
    preview: str
    page_count: int
    char_count: int
    parser_used: str   # "pdf"|"pdf+ocr"|"docx"|"xlsx"|"pptx"|"text"|"image-ocr"|"none"
    truncated: bool
    note: str | None = None


def _finalize(full_text: str, *, page_count: int, parser_used: str, note: str | None = None) -> ParsedDocument:
    preview = full_text[:PREVIEW_CHAR_LIMIT]
    return ParsedDocument(
        full_text=full_text,
        preview=preview,
        page_count=page_count,
        char_count=len(full_text),
        parser_used=parser_used,
        truncated=len(full_text) > PREVIEW_CHAR_LIMIT,
        note=note,
    )


def _ext(filename: str) -> str:
    dot = filename.rfind(".")
    return filename[dot:].lower() if dot != -1 else ""


async def parse_document(raw: bytes, mime: str, filename: str) -> ParsedDocument:
    mime = (mime or "").lower().split(";")[0].strip()
    ext = _ext(filename)

    if mime in _TEXT_MIMES or ext in _TEXT_EXTS:
        return _finalize(raw.decode("utf-8", errors="replace"), page_count=1, parser_used="text")
    if mime == "application/pdf" or ext == ".pdf":
        return await _parse_pdf(raw)
    if ext == ".docx" or "wordprocessingml" in mime:
        return _parse_docx(raw)
    if ext == ".xlsx" or "spreadsheetml" in mime:
        return _parse_xlsx(raw)
    if ext == ".pptx" or "presentationml" in mime:
        return _parse_pptx(raw)
    if mime in _IMAGE_MIMES or ext in {".png", ".jpg", ".jpeg", ".webp"}:
        return await _parse_image(raw, mime or f"image/{ext.lstrip('.')}")

    return _finalize("", page_count=0, parser_used="none", note=UNPARSEABLE_NOTE)


def _pdf_page_texts(raw: bytes) -> list[str]:
    """Return the extracted text of each PDF page (empty string when none)."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    return [(page.extract_text() or "").strip() for page in reader.pages]


async def _pdf_page_ocr(raw: bytes, page_index: int) -> str:
    """Render a single PDF page to a PNG and OCR it. Best-effort; "" on failure.

    Uses pypdfium2 (PDFium via pip wheels) so there is no system-binary dependency.
    """
    try:
        import pypdfium2 as pdfium

        pdf = pdfium.PdfDocument(raw)
        try:
            page = pdf[page_index]
            pil_image = page.render(scale=2.0).to_pil()  # ~144 DPI
        finally:
            pdf.close()
        buf = io.BytesIO()
        pil_image.save(buf, format="PNG")
        return await vision_ocr(buf.getvalue(), "image/png")
    except Exception:
        return ""


async def _parse_pdf(raw: bytes) -> ParsedDocument:
    try:
        texts = _pdf_page_texts(raw)
    except Exception:
        return _finalize("", page_count=0, parser_used="none", note="could not read PDF (corrupt or encrypted)")

    used_ocr = False
    out: list[str] = []
    for i, text in enumerate(texts):
        if text:
            out.append(text)
        else:
            ocr = await _pdf_page_ocr(raw, i)
            if ocr:
                used_ocr = True
                out.append(ocr)
    full = "\n\n".join(p for p in out if p)
    return _finalize(
        full,
        page_count=len(texts),
        parser_used="pdf+ocr" if used_ocr else "pdf",
        note=None if full else "no extractable text in PDF",
    )


def _parse_docx(raw: bytes) -> ParsedDocument:
    from docx import Document

    d = Document(io.BytesIO(raw))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return _finalize("\n".join(parts), page_count=1, parser_used="docx")


def _parse_xlsx(raw: bytes) -> ParsedDocument:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(",".join(cells))
    return _finalize("\n".join(parts), page_count=len(wb.worksheets), parser_used="xlsx")


def _parse_pptx(raw: bytes) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    slides = list(prs.slides)
    parts: list[str] = []
    for i, slide in enumerate(slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
    return _finalize("\n".join(parts), page_count=len(slides), parser_used="pptx")


async def _parse_image(raw: bytes, mime: str) -> ParsedDocument:
    text = await vision_ocr(raw, mime)
    return _finalize(
        text,
        page_count=1,
        parser_used="image-ocr",
        note=None if text else "no text extracted from image",
    )
